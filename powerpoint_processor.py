"""
powerpoint_processor.py
=======================
PowerPoint file operations and script execution.

FIXED:
- RgbColor import (was RGBColor)
- Added proper error handling
- Fixed type hints
- Added missing imports

Version: 1.0.1
"""

import traceback
from typing import Dict, Any, Optional, Tuple, List, Union
from enum import Enum
from datetime import datetime, date
from pathlib import Path
import warnings

# =============================================================================
# PPTX IMPORTS - FIXED: RgbColor not RGBColor
# =============================================================================

from pptx import Presentation
from pptx.util import Pt, Cm, Inches, Emu
from pptx.dml.color import RGBColor  # FIXED: lowercase 'gb'
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# =============================================================================
# ADI TEMPLATE IMPORTS
# =============================================================================

try:
    from adi_template_config import (
        COLORS, DIMENSIONS, TYPOGRAPHY, CONTAINERS, AMP_CONFIG,
        FOOTER_TEMPLATES, TITLE_CASE_LOWERCASE_WORDS, CLOSING_TAGLINE,
        COMPANY_URL, TABLE_CONFIG, ASSET_PATHS, VALIDATION,
        Confidentiality, SlideType, TableStyle, SectionSlideType, ChartColorScheme
    )
    
    from adi_pptx_generator import (
        ADIPresentation, ContentItem, ChartSeries, TableData,
        AssetManager, BrandValidator, AMPOverlayGenerator,
        to_title_case, format_bullet, hex_to_rgb, create_presentation
    )
    
    ADI_MODULES_AVAILABLE = True
except ImportError as e:
    ADI_MODULES_AVAILABLE = False
    warnings.warn(f"ADI template modules not available: {e}")
    
    # Define minimal fallbacks
    class Confidentiality(Enum):
        PUBLIC = "public"
        CONFIDENTIAL = "confidential"
        INTERNAL_ONLY = "internal_only"
    
    class TableStyle(Enum):
        RECOMMENDED_14PT = "recommended_14"
    
    class SectionSlideType(Enum):
        SECTION_TITLE = "section"
        KEY_MESSAGE = "key_message"


# =============================================================================
# BRAND COLORS
# =============================================================================

BRAND_COLORS = {
    "primary_blue": "0067B9",
    "dark_blue": "002855",
    "navy": "001A3D",
    "white": "FFFFFF",
    "light_gray": "F5F5F5",
    "text_dark": "333333",
    "text_light": "FFFFFF",
    "text_gray": "666666",
    "table_header": "B8D4E8",
    "table_alt_row": "E8F1F8",
}


# =============================================================================
# EXECUTION RESULT
# =============================================================================

class ExecutionResult:
    """Result of script execution."""
    
    def __init__(self, success: bool, presentation: Optional[Any],
                 error_message: str, traceback_str: str):
        self.success = success
        self.presentation = presentation
        self.error_message = error_message
        self.traceback = traceback_str
    
    @property
    def has_error(self) -> bool:
        return not self.success
    
    @property
    def workbook(self):
        """Alias for compatibility."""
        return self.presentation


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def create_content_items(items: List[Union[Tuple[str, int], Any]]) -> List:
    """Convert mixed content list to ContentItem objects."""
    if not ADI_MODULES_AVAILABLE:
        return items
    
    result = []
    for item in items:
        if isinstance(item, ContentItem):
            result.append(item)
        elif isinstance(item, tuple) and len(item) >= 2:
            text, level = item[0], item[1]
            extra_space = item[2] if len(item) > 2 else False
            result.append(ContentItem(text=text, level=level, extra_space=extra_space))
        else:
            result.append(ContentItem(text=str(item), level=1))
    return result


def create_table_data(headers: List[str], rows: List[List[str]]):
    """Create TableData from headers and rows."""
    if not ADI_MODULES_AVAILABLE:
        return {"headers": headers, "rows": rows}
    return TableData(headers=headers, rows=rows)


def create_chart_series(name: str, values: List[float]):
    """Create ChartSeries from name and values."""
    if not ADI_MODULES_AVAILABLE:
        return {"name": name, "values": values}
    return ChartSeries(name=name, values=values)


def validate_script(script: str) -> Tuple[bool, str]:
    """Validate script before execution."""
    if not script or not script.strip():
        return False, "Script is empty"
    
    dangerous_patterns = [
        ("import os", "Direct os import not allowed"),
        ("import sys", "Direct sys import not allowed"),
        ("import subprocess", "subprocess not allowed"),
        ("__import__", "Dynamic imports not allowed"),
        ("eval(", "eval() not allowed"),
        ("exec(", "Nested exec() not allowed"),
        ("open(", "Direct file operations not allowed"),
    ]
    
    for pattern, message in dangerous_patterns:
        if pattern in script:
            return False, message
    
    if "deck" not in script and "ADIPresentation" not in script:
        return False, "Script must create 'deck' variable using ADIPresentation"
    
    if ".save(" in script:
        return False, "Do not include save() in script - application handles saving"
    
    return True, ""


# =============================================================================
# SCRIPT EXECUTOR
# =============================================================================

class PPTXScriptExecutor:
    """Executes Python scripts that generate PowerPoint presentations."""
    
    def __init__(self):
        self._namespace = self._build_namespace()
    
    def _build_namespace(self) -> Dict[str, Any]:
        """Build the execution namespace."""
        namespace = {
            # Core pptx
            "Presentation": Presentation,
            "Pt": Pt,
            "Cm": Cm,
            "Inches": Inches,
            "Emu": Emu,
            "RgbColor": RGBColor,  # FIXED: correct name
            
            # Enums
            "PP_ALIGN": PP_ALIGN,
            "MSO_ANCHOR": MSO_ANCHOR,
            "MSO_SHAPE": MSO_SHAPE,
            "XL_CHART_TYPE": XL_CHART_TYPE,
            "XL_LEGEND_POSITION": XL_LEGEND_POSITION,
            
            # Chart
            "CategoryChartData": CategoryChartData,
            
            # Helpers
            "create_content_items": create_content_items,
            "create_table_data": create_table_data,
            "create_chart_series": create_chart_series,
            
            # Brand colors
            "BRAND_COLORS": BRAND_COLORS,
            
            # Date utilities
            "datetime": datetime,
            "date": date,
            "Path": Path,
        }
        
        # Add ADI modules if available
        if ADI_MODULES_AVAILABLE:
            namespace.update({
                "ADIPresentation": ADIPresentation,
                "ContentItem": ContentItem,
                "ChartSeries": ChartSeries,
                "TableData": TableData,
                "BrandValidator": BrandValidator,
                "AssetManager": AssetManager,
                "create_presentation": create_presentation,
                "to_title_case": to_title_case,
                "format_bullet": format_bullet,
                "hex_to_rgb": hex_to_rgb,
                
                # Configuration
                "COLORS": COLORS,
                "DIMENSIONS": DIMENSIONS,
                "TYPOGRAPHY": TYPOGRAPHY,
                "CONTAINERS": CONTAINERS,
                "AMP_CONFIG": AMP_CONFIG,
                "TABLE_CONFIG": TABLE_CONFIG,
                "VALIDATION": VALIDATION,
                
                # Enums
                "Confidentiality": Confidentiality,
                "SlideType": SlideType,
                "TableStyle": TableStyle,
                "SectionSlideType": SectionSlideType,
                "ChartColorScheme": ChartColorScheme,
                
                # Constants
                "FOOTER_TEMPLATES": FOOTER_TEMPLATES,
                "CLOSING_TAGLINE": CLOSING_TAGLINE,
                "COMPANY_URL": COMPANY_URL,
            })
        
        return namespace


    def validate_script(self, script: str) -> Tuple[bool, str]:
        """Validate script before execution."""
        if not script or not script.strip():
            return False, "Script is empty"
        
        dangerous_patterns = [
            ("import os", "Direct os import not allowed"),
            ("import sys", "Direct sys import not allowed"),
            ("import subprocess", "subprocess not allowed"),
            ("__import__", "Dynamic imports not allowed"),
            ("eval(", "eval() not allowed"),
            ("exec(", "Nested exec() not allowed"),
            ("open(", "Direct file operations not allowed"),
        ]
        
        for pattern, message in dangerous_patterns:
            if pattern in script:
                return False, message
        
        if "deck" not in script and "ADIPresentation" not in script:
            return False, "Script must create 'deck' variable using ADIPresentation"
        
        if ".save(" in script:
            return False, "Do not include save() in script - application handles saving"
        
        # Check for speaker_notes (warning, not error)
        if "speaker_notes" not in script:
            # This is a soft warning - script will still execute
            pass
        
        return True, ""


    def execute(self, script: str, validate: bool = True) -> ExecutionResult:
        """Execute script and return generated presentation."""
        if validate:
            is_valid, error_message = validate_script(script)
            if not is_valid:
                return ExecutionResult(False, None, error_message, "")
        
        namespace = self._namespace.copy()
        
        try:
            exec(script, namespace, namespace)
            
            if "deck" in namespace:
                presentation = namespace["deck"]
                if ADI_MODULES_AVAILABLE and isinstance(presentation, ADIPresentation):
                    return ExecutionResult(True, presentation, "", "")
                elif hasattr(presentation, 'prs'):  # Duck typing check
                    return ExecutionResult(True, presentation, "", "")
                else:
                    return ExecutionResult(False, None,
                        "'deck' is not a valid presentation object", "")
            else:
                return ExecutionResult(False, None,
                    "Script did not create 'deck' variable", "")
                
        except SyntaxError as e:
            return ExecutionResult(False, None,
                f"Syntax error at line {e.lineno}: {e.msg}",
                traceback.format_exc())
        except NameError as e:
            return ExecutionResult(False, None,
                f"Name error: {str(e)}",
                traceback.format_exc())
        except Exception as e:
            return ExecutionResult(False, None, str(e),
                traceback.format_exc())
    
    def save_presentation(self, presentation: Any, 
                          filepath: str, validate: bool = True) -> Tuple[bool, str]:
        """Save presentation to file."""
        try:
            if not filepath.endswith('.pptx'):
                filepath += '.pptx'
            
            if hasattr(presentation, 'save'):
                presentation.save(filepath, validate=validate)
            elif hasattr(presentation, 'prs'):
                presentation.prs.save(filepath)
            else:
                return False, "Invalid presentation object"
            
            return True, ""
        except Exception as e:
            return False, str(e)
    
    def get_namespace_info(self) -> Dict[str, List[str]]:
        """Get available namespace items."""
        info = {"classes": [], "functions": [], "constants": [], "enums": []}
        
        for name, obj in self._namespace.items():
            if isinstance(obj, type):
                if issubclass(obj, Enum) if isinstance(obj, type) else False:
                    info["enums"].append(name)
                else:
                    info["classes"].append(name)
            elif callable(obj):
                info["functions"].append(name)
            else:
                info["constants"].append(name)
        
        return info


# =============================================================================
# PRESENTATION BUILDER
# =============================================================================

class PresentationBuilder:
    """High-level builder for creating ADI presentations."""
    
    def __init__(self, 
                 confidentiality: Confidentiality = Confidentiality.PUBLIC,
                 year: Optional[int] = None,
                 assets_path: str = "assets"):
        if not ADI_MODULES_AVAILABLE:
            raise RuntimeError("ADI template modules required for PresentationBuilder")
        
        self.deck = ADIPresentation(
            confidentiality=confidentiality,
            year=year,
            assets_path=assets_path
        )
    
    def add_cover(self, title: str, subtitle: str = None,
                  background_image: str = None) -> 'PresentationBuilder':
        self.deck.add_cover_slide(title=title, subtitle=subtitle, background_image=background_image)
        return self
    
    def add_section(self, title: str, is_key_message: bool = False) -> 'PresentationBuilder':
        slide_type = SectionSlideType.KEY_MESSAGE if is_key_message else SectionSlideType.SECTION_TITLE
        self.deck.add_section_slide(title=title, slide_type=slide_type)
        return self
    
    def add_content(self, title: str, 
                    content: List[Union[Tuple[str, int], Any]],
                    dark_background: bool = False) -> 'PresentationBuilder':
        items = create_content_items(content)
        self.deck.add_content_slide(title=title, content=items, dark_background=dark_background)
        return self
    
    def add_two_column(self, title: str,
                       left_content: List,
                       right_content: List,
                       dark_background: bool = False) -> 'PresentationBuilder':
        left_items = create_content_items(left_content)
        right_items = create_content_items(right_content)
        self.deck.add_two_column_slide(
            title=title,
            left_content=left_items,
            right_content=right_items,
            dark_background=dark_background
        )
        return self
    
    def add_table(self, title: str, headers: List[str], rows: List[List[str]],
                  style: TableStyle = TableStyle.RECOMMENDED_14PT) -> 'PresentationBuilder':
        self.deck.add_table_slide(
            title=title,
            table_data=TableData(headers=headers, rows=rows),
            style=style
        )
        return self
    
    def add_chart(self, title: str, chart_type: str,
                  categories: List[str], series: List,
                  dark_background: bool = False) -> 'PresentationBuilder':
        self.deck.add_chart_slide(
            title=title,
            chart_type=chart_type,
            categories=categories,
            series=series,
            dark_background=dark_background
        )
        return self
    
    def add_closing(self) -> 'PresentationBuilder':
        self.deck.add_closing_slide()
        return self
    
    def build(self):
        return self.deck
    
    def save(self, filepath: str, validate: bool = True) -> str:
        return self.deck.save(filepath, validate=validate)


# =============================================================================
# MODULE INSTANCE
# =============================================================================

pptx_executor = PPTXScriptExecutor()


# =============================================================================
# CONVENIENCE FUNCTIONS
# =============================================================================

def execute_pptx_script(script: str) -> ExecutionResult:
    """Execute a PowerPoint generation script."""
    return pptx_executor.execute(script)


def save_pptx(presentation: Any, filepath: str) -> Tuple[bool, str]:
    """Save presentation to file."""
    return pptx_executor.save_presentation(presentation, filepath)


def get_pptx_namespace_info() -> Dict[str, List[str]]:
    """Get available namespace items."""
    return pptx_executor.get_namespace_info()


def quick_presentation(
    title: str,
    confidentiality: str = "public",
    slides: Optional[List[Dict[str, Any]]] = None
) -> Any:
    """Quick presentation creation from simple data."""
    if not ADI_MODULES_AVAILABLE:
        raise RuntimeError("ADI template modules required")
    
    conf_map = {
        "public": Confidentiality.PUBLIC,
        "confidential": Confidentiality.CONFIDENTIAL,
        "internal_only": Confidentiality.INTERNAL_ONLY,
    }
    
    builder = PresentationBuilder(
        confidentiality=conf_map.get(confidentiality.lower(), Confidentiality.PUBLIC)
    )
    
    builder.add_cover(title)
    
    if slides:
        for slide in slides:
            slide_type = slide.get("type", "content")
            
            if slide_type == "section":
                builder.add_section(
                    slide.get("title", "Section"),
                    slide.get("is_key_message", False)
                )
            elif slide_type == "content":
                builder.add_content(
                    slide.get("title", "Content"),
                    slide.get("content", []),
                    slide.get("dark_background", False)
                )
            elif slide_type == "two_column":
                builder.add_two_column(
                    slide.get("title", "Comparison"),
                    slide.get("left_content", []),
                    slide.get("right_content", [])
                )
            elif slide_type == "table":
                builder.add_table(
                    slide.get("title", "Table"),
                    slide.get("headers", []),
                    slide.get("rows", [])
                )
            elif slide_type == "chart":
                series_data = slide.get("series", [])
                series = [ChartSeries(s["name"], s["values"]) for s in series_data]
                builder.add_chart(
                    slide.get("title", "Chart"),
                    slide.get("chart_type", "column"),
                    slide.get("categories", []),
                    series
                )
            elif slide_type == "closing":
                builder.add_closing()
    else:
        builder.add_closing()
    
    return builder.build()
