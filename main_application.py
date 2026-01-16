"""
main_application.py
Main application orchestration for Analog Garage Workbench.
Entry point that ties all modules together - Enhanced with multi-output support.
Version 3.0
"""

import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
from datetime import datetime
import os
import traceback
from typing import Callable, Optional, Dict, Any, List, Tuple

from excel_importer import excel_importer

# Import configuration
from config import (
    APP_NAME, APP_VERSION, ADI_COLORS, UI_CONFIG, FEATURES, FORM_OPTIONS
)

# Import services
from services import prompt_service, context_service
from excel_processor import excel_executor
from session_manager import SessionManager


from powerpoint_processor import (
    execute_pptx_script,
    save_pptx,
    get_pptx_namespace_info,
    quick_presentation,
    PresentationBuilder,
    ExecutionResult,
    # ADI template classes
    ADIPresentation,
    ContentItem,
    ChartSeries,
    TableData,
    Confidentiality,
    TableStyle,
    SectionSlideType
)

# Import GUI components
from gui_components import (
    HeaderComponent,
    StatusBarComponent,
    InnovationFormComponent,
    PromptOutputComponent,
    ScriptInputComponent,
    ErrorDialogComponent,
    OutputSpecificInputsComponent
)

# Import templates
from templates import template_registry, build_enhanced_prompt


# =============================================================================
# OUTPUT TYPES CONFIGURATION (Embedded for self-contained operation)
# =============================================================================

OUTPUT_TYPES = {
    "excel_value_model": {
        "id": "excel_value_model",
        "name": "Excel Value Creation Model",
        "short_name": "Excel Model",
        "description": "Comprehensive Excel financial model with charts, projections, and sensitivity analysis",
        "output_format": "xlsx",
        "icon": "📊",
        "template_id": "value_creation_v3",
        "script_type": "excel",
        "color": "#28A745",
    },
    "powerpoint_pitch": {
        "id": "powerpoint_pitch",
        "name": "Executive Pitch Deck",
        "short_name": "PowerPoint Deck",
        "description": "Executive-ready PowerPoint presentation with key insights and visualizations",
        "output_format": "pptx",
        "icon": "📽️",
        "template_id": "executive_pitch_v1",
        "script_type": "powerpoint",
        "color": "#FF6600",
    },
    "word_gonogo": {
        "id": "word_gonogo",
        "name": "Deep Dive Summary - Recommend to Explore",
        "short_name": "Deep Dive Summary",
        "description": "Comprehensive Deep Dive assessment with exploration recommendation for Analog Garage",
        "output_format": "docx",
        "icon": "📋",
        "template_id": "gonogo_report_v1",
        "script_type": "word",
        "color": "#0067B9",
    },
}

# List for dropdown
OUTPUT_TYPE_OPTIONS = [
    f"{v['icon']} {v['name']}" for v in OUTPUT_TYPES.values()
]

# Mapping from display name to ID
OUTPUT_TYPE_DISPLAY_TO_ID = {
    f"{v['icon']} {v['name']}": k for k, v in OUTPUT_TYPES.items()
}

# Script instructions by output type
SCRIPT_INSTRUCTIONS = {
    "excel": {
        "title": "Excel Script Input",
        "instruction": "💡 Paste the AI-generated Python script below. The script must create a workbook object named 'wb' using openpyxl. Do NOT include wb.save().",
        "button_text": "▶️ Execute & Save Excel",
        "file_extension": ".xlsx",
        "file_type": "Excel files",
    },
    "powerpoint": {
        "title": "PowerPoint Script Input",
        "instruction": "💡 Paste the AI-generated Python script below. The script must create a presentation object named 'prs' using python-pptx. Do NOT include prs.save().",
        "button_text": "▶️ Execute & Save PowerPoint",
        "file_extension": ".pptx",
        "file_type": "PowerPoint files",
    },
    "word": {
        "title": "Word Document Script Input",
        "instruction": "💡 Paste the AI-generated Python script below. The script must create a document object named 'doc' using python-docx. Do NOT include doc.save().",
        "button_text": "▶️ Execute & Save Word Doc",
        "file_extension": ".docx",
        "file_type": "Word files",
    },
}


# =============================================================================
# OUTPUT TYPE SELECTOR COMPONENT
# =============================================================================

class OutputTypeSelectorComponent(tk.Frame):
    """Component for selecting the output type (Excel/PowerPoint/Word)."""
    
    def __init__(self, parent, on_change: Optional[Callable] = None):
        super().__init__(parent, bg=ADI_COLORS["white"])
        self.on_change = on_change
        self._create_widgets()
    
    def _create_widgets(self):
        # Container with visual emphasis
        container = tk.Frame(self, bg=ADI_COLORS["light_blue"], padx=15, pady=12)
        container.pack(fill=tk.X, padx=10, pady=10)
        
        # Title row
        title_frame = tk.Frame(container, bg=ADI_COLORS["light_blue"])
        title_frame.pack(fill=tk.X)
        
        tk.Label(
            title_frame,
            text="🎯 Select Output Type",
            font=(UI_CONFIG["font_family"], 12, "bold"),
            bg=ADI_COLORS["light_blue"],
            fg=ADI_COLORS["dark_blue"]
        ).pack(side=tk.LEFT)
        
        # Description
        tk.Label(
            container,
            text="Choose the type of deliverable you want to generate:",
            font=(UI_CONFIG["font_family"], 10),
            bg=ADI_COLORS["light_blue"],
            fg=ADI_COLORS["dark_gray"]
        ).pack(anchor="w", pady=(5, 10))
        
        # Dropdown
        dropdown_frame = tk.Frame(container, bg=ADI_COLORS["light_blue"])
        dropdown_frame.pack(fill=tk.X)
        
        self.output_type_combo = ttk.Combobox(
            dropdown_frame,
            width=45,
            state="readonly",
            values=OUTPUT_TYPE_OPTIONS,
            font=(UI_CONFIG["font_family"], 11)
        )
        
        # Set default value
        if OUTPUT_TYPE_OPTIONS:
            self.output_type_combo.set(OUTPUT_TYPE_OPTIONS[0])
        
        self.output_type_combo.pack(side=tk.LEFT)
        
        # Bind change event
        self.output_type_combo.bind("<<ComboboxSelected>>", self._on_selection_change)
        
        # Description label (updates based on selection)
        self.desc_label = tk.Label(
            container,
            text=self._get_description(),
            font=(UI_CONFIG["font_family"], 9, "italic"),
            bg=ADI_COLORS["light_blue"],
            fg=ADI_COLORS["text_gray"],
            wraplength=500,
            justify=tk.LEFT
        )
        self.desc_label.pack(anchor="w", pady=(8, 0))
    
    def _on_selection_change(self, event=None):
        """Handle selection change."""
        self.desc_label.config(text=self._get_description())
        if self.on_change:
            self.on_change(self.get_output_type_id())
    
    def _get_description(self) -> str:
        """Get description for current selection."""
        output_id = self.get_output_type_id()
        if output_id and output_id in OUTPUT_TYPES:
            return OUTPUT_TYPES[output_id]["description"]
        return ""
    
    def get_output_type_id(self) -> str:
        """Get the selected output type ID."""
        display_name = self.output_type_combo.get()
        return OUTPUT_TYPE_DISPLAY_TO_ID.get(display_name, "excel_value_model")
    
    def get_output_type_config(self) -> dict:
        """Get full configuration for selected output type."""
        output_id = self.get_output_type_id()
        return OUTPUT_TYPES.get(output_id, OUTPUT_TYPES["excel_value_model"])
    
    def set_output_type(self, output_id: str):
        """Set the output type by ID."""
        for display_name, id_ in OUTPUT_TYPE_DISPLAY_TO_ID.items():
            if id_ == output_id:
                self.output_type_combo.set(display_name)
                self._on_selection_change()
                break


# =============================================================================
# ENHANCED SCRIPT INPUT COMPONENT
# =============================================================================

class ScriptInputComponentEnhanced(tk.Frame):
    """Enhanced script input component with dynamic instructions based on output type."""
    
    def __init__(self, parent, on_execute: Callable, on_clear: Callable):
        super().__init__(parent, bg=ADI_COLORS["white"])
        self.on_execute = on_execute
        self.on_clear = on_clear
        self._create_widgets()
    
    def _create_widgets(self):
        # Info frame with dynamic content
        self.info_frame = tk.Frame(self, bg=ADI_COLORS["light_blue"], pady=10)
        self.info_frame.pack(fill=tk.X, pady=(0, 10))
        
        self.instruction_label = tk.Label(
            self.info_frame,
            text=SCRIPT_INSTRUCTIONS["excel"]["instruction"],
            font=(UI_CONFIG["font_family"], UI_CONFIG["small_font_size"]),
            bg=ADI_COLORS["light_blue"],
            fg=ADI_COLORS["dark_blue"],
            wraplength=900
        )
        self.instruction_label.pack(padx=10)
        
        # Title label
        self.title_label = tk.Label(
            self,
            text="Python Script Input:",
            font=(UI_CONFIG["font_family"], UI_CONFIG["body_font_size"], "bold"),
            bg=ADI_COLORS["white"],
            fg=ADI_COLORS["dark_gray"]
        )
        self.title_label.pack(anchor="w", pady=(0, 5))
        
        # Script input area
        self.script_input = scrolledtext.ScrolledText(
            self,
            font=(UI_CONFIG["code_font_family"], UI_CONFIG["body_font_size"]),
            bg=ADI_COLORS["code_bg"],
            fg=ADI_COLORS["code_fg"],
            insertbackground=ADI_COLORS["accent_orange"],
            wrap=tk.WORD
        )
        self.script_input.pack(fill=tk.BOTH, expand=True)
        
        # Button frame
        btn_frame = tk.Frame(self, bg=ADI_COLORS["white"], pady=15)
        btn_frame.pack(fill=tk.X)
        
        tk.Button(
            btn_frame,
            text="🗑️ Clear",
            command=self.on_clear,
            bg=ADI_COLORS["medium_gray"],
            fg="black",
            font=(UI_CONFIG["font_family"], UI_CONFIG["small_font_size"]),
            width=10,
            height=2,
            cursor="hand2"
        ).pack(side=tk.LEFT)
        
        self.execute_btn = tk.Button(
            btn_frame,
            text=SCRIPT_INSTRUCTIONS["excel"]["button_text"],
            command=self.on_execute,
            bg=ADI_COLORS["accent_orange"],
            fg="black",
            font=(UI_CONFIG["font_family"], 13, "bold"),
            width=24,
            height=2,
            cursor="hand2"
        )
        self.execute_btn.pack(side=tk.RIGHT)
    
    def update_instructions(self, instruction: str, button_text: str):
        """Update the instruction text and button label."""
        self.instruction_label.config(text=instruction)
        self.execute_btn.config(text=button_text)
    
    def get_script(self) -> str:
        return self.script_input.get("1.0", tk.END).strip()
    
    def clear(self):
        self.script_input.delete("1.0", tk.END)


# =============================================================================
# MAIN APPLICATION CLASS
# =============================================================================

class AnalogGarageWorkbench:
    """
    Main application class that orchestrates all components.
    Enhanced with multi-output type support (Excel, PowerPoint, Word).
    """
    
    def __init__(self):
        # Initialize main window
        self.root = tk.Tk()
        self.root.title(f"{APP_NAME} v3.0")
        self.root.geometry(f"{UI_CONFIG['window_width']}x{UI_CONFIG['window_height']}")
        self.root.configure(bg=ADI_COLORS["light_gray"])
        
        # Initialize session manager
        self.session = SessionManager()
        
        # Track current output type
        self.current_output_type = "excel_value_model"
        
        # Build UI
        self._create_menu()
        self._create_header()
        self._create_notebook()
        self._create_status_bar()
        
        # Bind keyboard shortcuts
        self._bind_shortcuts()
        
        self.session.info("Application initialized with multi-output support")
    
    # =========================================================================
    # UI CREATION
    # =========================================================================
    
    def _create_menu(self):
        """Create application menu bar."""
        menubar = tk.Menu(self.root)
        self.root.config(menu=menubar)
        
        # File menu
        file_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="File", menu=file_menu)
        file_menu.add_command(label="Export Log...", command=self._export_log)
        file_menu.add_separator()
        file_menu.add_command(label="Exit", command=self.root.quit)
        
        # Output Type menu
        output_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="Output Type", menu=output_menu)
        for output_id, config in OUTPUT_TYPES.items():
            output_menu.add_command(
                label=f"{config['icon']} {config['name']}",
                command=lambda oid=output_id: self._select_output_type(oid)
            )
        
        # Templates menu
        templates_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="Templates", menu=templates_menu)
        for template in template_registry.get_all():
            templates_menu.add_command(
                label=template.name,
                command=lambda t=template.id: self._select_template(t)
            )
        
        # Help menu
        help_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="Help", menu=help_menu)
        help_menu.add_command(label="User Guide", command=self._show_help)
        help_menu.add_command(label="About", command=self._show_about)
    
    def _create_header(self):
        """Create application header."""
        self.header = HeaderComponent(self.root, demo_mode=FEATURES["demo_mode"])
        self.header.pack(fill=tk.X)
    
    def _create_notebook(self):
        """Create tabbed interface."""
        style = ttk.Style()
        style.configure('TNotebook.Tab', font=('Arial', 11, 'bold'), padding=[15, 8])
        
        self.notebook = ttk.Notebook(self.root)
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=15, pady=10)
        
        # Tab 1: Prompt Generator
        self._create_prompt_tab()
        
        # Tab 2: Script Executor
        self._create_script_tab()
    
    def _create_prompt_tab(self):
        """Create the Prompt Generator tab with output type selector and dynamic forms."""
        prompt_frame = tk.Frame(self.notebook, bg=ADI_COLORS["white"])
        self.notebook.add(prompt_frame, text="  🎯 1. Generate Prompt  ")
        
        # Create a PanedWindow for resizable sections
        paned = ttk.PanedWindow(prompt_frame, orient=tk.VERTICAL)
        paned.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # TOP SECTION: Output selector + Base form
        top_frame = tk.Frame(paned, bg=ADI_COLORS["white"])
        paned.add(top_frame, weight=1)
        
        # Output Type Selector
        self.output_selector = OutputTypeSelectorComponent(
            top_frame,
            on_change=self._on_output_type_change
        )
        self.output_selector.pack(fill=tk.X, padx=5, pady=(10, 5))
        
        # Base Innovation Form (scrollable)
        self.form = InnovationFormComponent(
            top_frame,
            on_generate=self._generate_prompt,
            on_clear=self._clear_form,
            on_import=self._import_from_excel
        )
        self.form.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # MIDDLE SECTION: Output-specific inputs
        middle_frame = tk.Frame(paned, bg=ADI_COLORS["white"])
        paned.add(middle_frame, weight=1)
        
        # Collapsible output-specific section
        self.specific_inputs_visible = tk.BooleanVar(value=True)
        
        toggle_frame = tk.Frame(middle_frame, bg=ADI_COLORS["white"])
        toggle_frame.pack(fill=tk.X, padx=10, pady=5)
        
        self.specific_toggle_btn = tk.Button(
            toggle_frame,
            text="▼ Output-Specific Inputs (Recommended)",
            command=self._toggle_specific_inputs,
            bg=ADI_COLORS["white"],
            fg=ADI_COLORS["primary_blue"],
            font=(UI_CONFIG["font_family"], 10, "bold"),
            relief=tk.FLAT,
            cursor="hand2"
        )
        self.specific_toggle_btn.pack(anchor="w")
        
        # Output-specific inputs component
        self.output_specific_inputs = OutputSpecificInputsComponent(
            middle_frame,
            bg_color=ADI_COLORS["white"]
        )
        self.output_specific_inputs.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # Connect the form to the output-specific component for export/import
        self.form.set_output_specific_component(self.output_specific_inputs)


        # BOTTOM SECTION: Generated prompt output
        bottom_frame = tk.Frame(paned, bg=ADI_COLORS["white"])
        paned.add(bottom_frame, weight=1)
        
        # Separator
        ttk.Separator(bottom_frame, orient="horizontal").pack(fill=tk.X, padx=15, pady=5)
        
        # Output component
        self.prompt_output = PromptOutputComponent(
            bottom_frame,
            on_copy=self._copy_prompt,
            on_export=self._export_prompt,
            on_generate=self._generate_prompt
        )
        self.prompt_output.pack(fill=tk.BOTH, expand=True, padx=15, pady=(5, 15))

    def _toggle_specific_inputs(self):
        """Toggle visibility of output-specific inputs."""
        if self.specific_inputs_visible.get():
            self.output_specific_inputs.pack_forget()
            self.specific_toggle_btn.config(text="▶ Output-Specific Inputs (Recommended)")
            self.specific_inputs_visible.set(False)
        else:
            self.output_specific_inputs.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            self.specific_toggle_btn.config(text="▼ Output-Specific Inputs (Recommended)")
            self.specific_inputs_visible.set(True)

    def _on_output_type_change(self, output_type_id: str):
        """Handle output type change - update script tab, specific inputs, and form."""
        self.current_output_type = output_type_id
        config = OUTPUT_TYPES.get(output_type_id, OUTPUT_TYPES["excel_value_model"])
        
        # Update script tab
        self._update_script_tab_for_output_type(config)
        
        # Update output-specific inputs
        self.output_specific_inputs.set_output_type(output_type_id)
        
        # Sync output type to form for template generation
        self.form.current_output_type = output_type_id
        
        # Log the change
        self.session.info(f"Output type changed to: {config['name']}")
        self._update_status(f"Output: {config['icon']} {config['short_name']}", config['color'])


    
    def _update_script_tab_for_output_type(self, config: dict):
        """Update the script tab UI based on selected output type."""
        script_type = config.get("script_type", "excel")
        instructions = SCRIPT_INSTRUCTIONS.get(script_type, SCRIPT_INSTRUCTIONS["excel"])
        
        # Update script input component
        if hasattr(self, 'script_input'):
            self.script_input.update_instructions(
                instructions["instruction"],
                instructions["button_text"]
            )
    
    def _import_from_excel(self, filepath: str):
        """Handle importing data from Excel file including output-specific fields."""
        self.session.info(f"Importing from: {filepath}")
        
        result = excel_importer.import_from_excel(filepath)
        
        if result.success:
            # Populate base form fields
            self.form.populate_from_data(result.data)
            
            # Set output type if detected and different from current
            if result.output_type and result.output_type != self.current_output_type:
                self._select_output_type(result.output_type)
                # Small delay to allow UI to update
                self.root.update_idletasks()
            
            # Populate output-specific fields
            if result.specific_data:
                self.form.populate_specific_fields(result.specific_data)
            
            # Count fields
            base_field_count = len([v for v in result.data.values() if v])
            specific_field_count = len([v for v in result.specific_data.values() if v])
            total_count = base_field_count + specific_field_count
            
            message = f"Successfully imported {total_count} fields from:\n{os.path.basename(filepath)}"
            if specific_field_count > 0:
                message += f"\n\n• {base_field_count} base fields\n• {specific_field_count} output-specific fields"
            if result.output_type:
                message += f"\n\nOutput type: {result.output_type}"
            if result.warnings:
                message += f"\n\nWarnings:\n• " + "\n• ".join(result.warnings[:5])
            
            self.session.info(f"Import successful: {total_count} fields ({base_field_count} base, {specific_field_count} specific)")
            self._update_status(f"✓ Imported from {os.path.basename(filepath)}", ADI_COLORS["success_green"])
            messagebox.showinfo("Import Successful", message)
        else:
            error_msg = "\n".join(result.errors)
            self.session.error(f"Import failed: {error_msg}")
            messagebox.showerror("Import Failed", error_msg)

    
    def _create_script_tab(self):
        """Create the Script Executor tab with dynamic instructions."""
        script_frame = tk.Frame(self.notebook, bg=ADI_COLORS["white"])
        self.notebook.add(script_frame, text="  ⚡ 2. Execute Script  ")
        
        self.script_input = ScriptInputComponentEnhanced(
            script_frame,
            on_execute=self._execute_script,
            on_clear=self._clear_script
        )
        self.script_input.pack(fill=tk.BOTH, expand=True, padx=15, pady=15)
    
    def _create_status_bar(self):
        """Create status bar."""
        self.status_bar = StatusBarComponent(self.root, version="3.0")
        self.status_bar.pack(fill=tk.X, side=tk.BOTTOM)
    
    def _bind_shortcuts(self):
        """Bind keyboard shortcuts."""
        self.root.bind("<Command-q>", lambda e: self.root.quit())
        self.root.bind("<Control-q>", lambda e: self.root.quit())
    
    # =========================================================================
    # PROMPT GENERATION - Updated for multiple output types
    # =========================================================================
    
    def _generate_prompt(self):
        """Generate prompt based on selected output type with specific inputs."""
        values = self.form.get_values()
        
        # Get output-specific inputs
        specific_values = self.output_specific_inputs.get_values()
        
        # Get current output type configuration
        output_config = self.output_selector.get_output_type_config()
        template_id = output_config.get("template_id", "value_creation_v3")
        
        # Build context with both base and specific values
        context = self._build_context_from_values(values)
        context["specific_inputs"] = specific_values
        
        # Validate required fields
        if not context["innovation_name"]:
            messagebox.showwarning("Required Field", "Please enter an Innovation Name.")
            return
        if not context["industry"]:
            messagebox.showwarning("Required Field", "Please enter a Target Industry.")
            return
        if not context["innovation_description"]:
            messagebox.showwarning("Required Field", "Please enter an Innovation Description.")
            return
        
        # Generate prompt based on output type
        prompt = self._generate_prompt_for_type(context, template_id, output_config)
        
        # Update UI
        self.prompt_output.set_content(prompt)
        context_service.set_current(context)
        self.session.set_prompt(prompt)
        self.session.set_innovation_context(context)
        
        output_name = output_config.get("short_name", "Prompt")
        self._update_status(
            f"✓ {output_name} prompt generated for: {values['innovation_name']}", 
            output_config.get("color", ADI_COLORS["success_green"])
        )
        
        # Show confirmation
        self._show_generation_confirmation(output_config, values)

    
    def _build_context_from_values(self, values: dict) -> dict:
        """Build context dictionary from form values."""
        # Build conditional sections
        problem_section = ""
        if values.get("problem_statement"):
            problem_section = f"**Problem Statement:**\n{values['problem_statement']}"
        
        customer_section = ""
        if values.get("target_customer"):
            customer_section = f"**Target Customer:** {values['target_customer']}"
        
        advantage_section = ""
        if values.get("competitive_advantage"):
            advantage_section = f"**Competitive Advantage:** {values['competitive_advantage']}"
        
        # Market section
        market_items = []
        if values.get("tam"):
            market_items.append(f"• Total Addressable Market: {values['tam']}")
        if values.get("target_penetration"):
            market_items.append(f"• Target Market Penetration: {values['target_penetration']}")
        if values.get("price_point"):
            market_items.append(f"• Target Price Point: {values['price_point']}")
        market_section = "\n".join(market_items) if market_items else "(No market estimates provided)"
        
        # Regulatory section
        regulatory_items = []
        if values.get("regulatory_pathway"):
            regulatory_items.append(f"• Regulatory Pathway: {values['regulatory_pathway']}")
        if values.get("ip_status"):
            regulatory_items.append(f"• IP Status: {values['ip_status']}")
        regulatory_section = "\n".join(regulatory_items) if regulatory_items else "(No regulatory/IP context provided)"
        
        # Risks section
        risks_section = ""
        if values.get("key_risks"):
            risks_section = f"**Key Risks Identified:**\n{values['key_risks']}"
        else:
            risks_section = "(No specific risks identified)"
        
        # Assumptions section
        assumptions_section = ""
        if values.get("key_assumptions"):
            assumptions_section = f"**User-Provided Assumptions:**\n{values['key_assumptions']}"
        else:
            assumptions_section = "(No specific assumptions provided)"
        
        return {
            "innovation_name": values.get("innovation_name", ""),
            "innovation_description": values.get("innovation_description", ""),
            "industry": values.get("industry", ""),
            "geographic_scope": values.get("geographic_scope", "Global"),
            "analysis_timeframe": values.get("analysis_timeframe", "Year 1 at Scale"),
            "innovation_stage": values.get("innovation_stage", "Concept"),
            "currency": values.get("currency", "USD"),
            "problem_section": problem_section,
            "customer_section": customer_section,
            "advantage_section": advantage_section,
            "market_section": market_section,
            "regulatory_section": regulatory_section,
            "risks_section": risks_section,
            "assumptions_section": assumptions_section,
            # Raw values for conditional logic
            "problem_statement": values.get("problem_statement", ""),
            "target_customer": values.get("target_customer", ""),
            "competitive_advantage": values.get("competitive_advantage", ""),
            "tam": values.get("tam", ""),
            "target_penetration": values.get("target_penetration", ""),
            "price_point": values.get("price_point", ""),
            "regulatory_pathway": values.get("regulatory_pathway", ""),
            "ip_status": values.get("ip_status", ""),
            "key_risks": values.get("key_risks", ""),
            "key_assumptions": values.get("key_assumptions", ""),
        }
    
    def _generate_prompt_for_type(self, context: dict, template_id: str, output_config: dict) -> str:
        """Generate prompt based on output type."""
        script_type = output_config.get("script_type", "excel")
        
        if script_type == "excel":
            # Use existing enhanced builder for Excel
            return build_enhanced_prompt(context)
        elif script_type == "powerpoint":
            # Use PowerPoint template
            return self._generate_powerpoint_prompt(context)
        elif script_type == "word":
            # Use Word/Deep Dive template
            return self._generate_word_prompt(context)
        else:
            # Fallback to template registry
            template = template_registry.get(template_id)
            if template:
                try:
                    return template.render(context)
                except Exception as e:
                    self.session.error(f"Template render error: {e}")
                    return f"Error rendering template: {e}"
            return build_enhanced_prompt(context)
    
    
    def _generate_powerpoint_prompt(self, context: dict) -> str:
        """
        Generate enhanced prompt for Executive Pitch Deck with:
        1. Complete ADI brand theme integration
        2. Full slide method documentation with visual guidance
        3. Embedded speaker notes requirements
        4. Aesthetic design principles
        """
        specific = context.get("specific_inputs", {})
        
        # Build specific input sections with verification instructions
        problem_guidance = ""
        if specific.get("problem_primary") or specific.get("problem_secondary"):
            problem_guidance = f"""
    ───────────────────────────────────────────────────────────────────────────────
    USER-PROVIDED PROBLEM STATEMENT — VERIFY & ENHANCE
    ───────────────────────────────────────────────────────────────────────────────

    **BASELINE INPUT (treat as authoritative starting point):**

    Primary Problem Statement:
    {specific.get('problem_primary', '(Not provided)')}

    Supporting Problems / Pain Points:
    {specific.get('problem_secondary', '(Not provided)')}

    Key Statistics & Evidence:
    {specific.get('problem_stats', '(Not provided)')}

    **YOUR TASK:**
    1. VERIFY: Cross-check statistics against your knowledge. Flag questionable claims.
    2. ENHANCE: Sharpen language for executive impact. Convert passive to active, urgent statements.
    3. QUANTIFY: Add specific numbers where you can confidently estimate.
    4. REFRAME: Express problems in business impact terms (revenue lost, time wasted, lives affected).
    """

        solution_guidance = ""
        if specific.get("solution_how") or specific.get("solution_features"):
            solution_guidance = f"""
    ───────────────────────────────────────────────────────────────────────────────
    USER-PROVIDED SOLUTION — VERIFY & ENHANCE
    ───────────────────────────────────────────────────────────────────────────────

    **BASELINE INPUT:**

    How It Works:
    {specific.get('solution_how', '(Not provided)')}

    Key Features & Capabilities:
    {specific.get('solution_features', '(Not provided)')}

    Quantified Benefits:
    {specific.get('solution_benefits', '(Not provided)')}

    **YOUR TASK:**
    1. VERIFY: Ensure claimed benefits are plausible given the technology.
    2. ENHANCE: Lead with OUTCOMES, not technology. Use active verbs.
    3. PRIORITIZE: Identify the single most compelling benefit.
    4. DIFFERENTIATE: Highlight what's genuinely unique.
    """

        value_guidance = ""
        if specific.get("value_prop") or specific.get("differentiators"):
            value_guidance = f"""
    ───────────────────────────────────────────────────────────────────────────────
    USER-PROVIDED VALUE PROPOSITION — VERIFY & SHARPEN
    ───────────────────────────────────────────────────────────────────────────────

    **BASELINE INPUT:**

    Core Value Proposition:
    {specific.get('value_prop', '(Not provided)')}

    Key Differentiators:
    {specific.get('differentiators', '(Not provided)')}

    **YOUR TASK:**
    1. VERIFY: Is this differentiation defensible?
    2. SHARPEN: Compress to ONE sentence that passes the "so what?" test.
    3. ANCHOR: Tie differentiation to specific, verifiable capabilities.
    """

        ask_guidance = ""
        if specific.get("the_ask") or specific.get("use_of_funds"):
            ask_guidance = f"""
    ───────────────────────────────────────────────────────────────────────────────
    USER-PROVIDED ASK — VERIFY & STRUCTURE
    ───────────────────────────────────────────────────────────────────────────────

    **BASELINE INPUT:**

    Investment/Resource Request:
    {specific.get('the_ask', '(Not provided)')}

    Use of Funds / Resources:
    {specific.get('use_of_funds', '(Not provided)')}

    **YOUR TASK:**
    1. VERIFY: Is the ask proportionate to the opportunity?
    2. STRUCTURE: Break into clear buckets (R&D, Clinical, Commercial).
    3. TIMELINE: Add expected milestones for each investment tranche.
    """

        return f"""You are an expert executive presentation strategist specializing in high-stakes pitch decks for the {context['industry']} sector.

    ═══════════════════════════════════════════════════════════════════════════════
    ADI BRAND THEME — MANDATORY VISUAL STANDARDS
    ═══════════════════════════════════════════════════════════════════════════════

    **CORPORATE COLOR PALETTE (Use these exact values):**

    | Role | Hex Code | Usage |
    |------|----------|-------|
    | Primary Navy | #00325C | Headers, key metrics, primary emphasis |
    | ADI Blue | #0067B9 | Secondary headers, hyperlinks, accents |
    | Sky Blue | #1B9CD0 | Charts, highlights, positive indicators |
    | Purple | #8637BA | Differentiation, innovation themes |
    | Green | #179963 | Success, growth, positive outcomes |
    | Yellow | #FED141 | Warnings, attention, caution |
    | Red | #C81A28 | Urgency, problems, negative indicators |
    | Dark Text | #000000 | Body text on light backgrounds |
    | Light Text | #FFFFFF | Text on dark backgrounds |
    | Gray | #9EA1AE | Secondary text, footnotes |

    **TYPOGRAPHY:**
    - Title Font: Barlow Medium (bold for emphasis)
    - Body Font: Barlow
    - Minimum title size: 28pt
    - Body text: 18pt (Level 1), 16pt (Level 2), 14pt (Level 3)

    **VISUAL DESIGN PRINCIPLES:**
    1. ONE dominant visual element per slide
    2. Maximum 5 bullet points per slide
    3. Use icons (✓, →, •, ★) to enhance scannability
    4. Dark backgrounds (#00325C) for urgency/problems
    5. Light backgrounds (#FFFFFF) for solutions/data
    6. Charts use accent colors in sequence: Navy → Sky Blue → Purple → Green

    ═══════════════════════════════════════════════════════════════════════════════
    INNOVATION OVERVIEW
    ═══════════════════════════════════════════════════════════════════════════════

    **Innovation Name:** {context['innovation_name']}
    **Target Industry:** {context['industry']}
    **Geographic Scope:** {context['geographic_scope']}
    **Stage:** {context['innovation_stage']}

    **DESCRIPTION:**
    {context['innovation_description']}

    {context.get('problem_section', '')}
    {context.get('customer_section', '')}
    {context.get('advantage_section', '')}

    **MARKET CONTEXT:**
    {context.get('market_section', '(No market data provided)')}

    {problem_guidance}
    {solution_guidance}
    {value_guidance}
    {ask_guidance}

    ═══════════════════════════════════════════════════════════════════════════════
    COMPLETE SLIDE METHOD REFERENCE — ADIPresentation CLASS
    ═══════════════════════════════════════════════════════════════════════════════

    Very Important: Use the following as guidance only.
    Add shape elements and other python-pptx elements to augment these base templates and use creative license on the flow of slides to maximise impact.

    **1. COVER SLIDE** — First impression, sets the tone
    ```python
    deck.add_cover_slide(
        title="Presentation Title",           # Large, bold, title case
        subtitle="Optional Subtitle",          # Smaller, below title
        background_image="path/to/image.jpg", # Optional custom background
        use_default_bg=True,                  # Use ADI branded background
        speaker_notes="Notes for presenter"
    )

    Visual: Dark blue background with AMP triangle overlay, white text, logo top-left
    2. SECTION SLIDE — Divider for major sections or key messages
    deck.add_section_slide(
        title="Section Title or Key Statement",
        slide_type=SectionSlideType.SECTION_TITLE,  # or KEY_MESSAGE for larger text
        background_image=None,                       # Optional custom background
        speaker_notes="Transition notes"
    )

    Visual: Full-bleed dark background, centered white text, AMP overlay
    3. CONTENT SLIDE — Standard bullet-point content
    deck.add_content_slide(
        title="Slide Title",
        content=[
            ContentItem("Main point (Level 1)", level=1),
            ContentItem("Supporting detail", level=2),
            ContentItem("Sub-detail", level=3),
            ContentItem("Point with extra spacing", level=1, extra_space=True),
            ContentItem("Highlighted point", level=1, highlight=True),
        ],
        dark_background=False,  # True for urgency/problems
        speaker_notes="Key talking points..."
    )

    Visual: Logo top-right, title at top, hierarchical bullets below
    4. METRIC HIGHLIGHT SLIDE — Single impactful statistic
    deck.add_metric_highlight_slide(
        title="The Scale of the Problem",
        metric_value="\$2.2B",                    # Large, centered number
        metric_label="Total Addressable Market", # Explanation below
        supporting_text="Source: Industry Report 2025",
        dark_background=True,
        speaker_notes="Emphasize the magnitude..."
    )

    Visual: Metric dominates center (72pt), label below (24pt), dark = urgency
    5. COMPARISON SLIDE — Before/After or Problem/Solutio   
    deck.add_comparison_slide(
        title="Current State vs. Our Solution",
        before_header="Today's Reality",
        before_items=[
            "Manual review required",
            "4-6 hour delays",
            "Specialist dependency"
        ],
        after_header="With Our Solution",
        after_items=[
            "Automated detection",
            "15-minute alerts",
            "Any clinician can act"
        ],
        before_color="#C81A28",  # Red for problems
        after_color="#179963",   # Green for solutions
        speaker_notes="Contrast is key..."
    )
    Visual: Two columns, ✗ icons for before (red), ✓ icons for after (green)
    6. TWO-COLUMN SLIDE — Side-by-side content
        deck.add_two_column_slide(
        title="Feature Comparison",
        left_header="Our Approach",
        left_content=[
            ContentItem("Advantage 1", level=1),
            ContentItem("Detail", level=2),
        ],
        right_header="Competitor",
        right_content=[
            ContentItem("Their approach", level=1),
            ContentItem("Limitation", level=2),
        ],
        dark_background=False,
        speaker_notes="Highlight key differences..."
    )
    Visual: Equal columns, optional headers in bold
    7. TABLE SLIDE — Structured data presentation   
    deck.add_table_slide(
        title="Financial Projections",
        table_data=TableData(
            headers=["Year", "Revenue", "Growth"],
            rows=[
                ["2026", "\$2M", "—"],
                ["2027", "\$8M", "300%"],
                ["2028", "\$25M", "212%"]
            ],
            highlight_row=2  # Zero-indexed, highlights row visually
        ),
        style=TableStyle.RECOMMENDED_14PT,
        speaker_notes="Focus on the highlighted row..."
    )
    Visual: Blue header row (#B8D4E8), alternating row colors, bold highlight
    8. CHART SLIDE — Data visualization
    deck.add_chart_slide(
        title="Market Opportunity",
        chart_type="column",  # Options: bar, column, line, pie, area, stacked_bar, stacked_column, stacked_area
        categories=["Q1", "Q2", "Q3", "Q4"],
        series=[
            ChartSeries("Revenue", [100, 150, 200, 280]),
            ChartSeries("Target", [120, 160, 210, 300])
        ],
        dark_background=False,
        chart_title="Revenue vs Target",  # Optional chart title
        speaker_notes="Trend shows acceleration..."
    )
    Visual: Chart uses accent colors automatically, legend at bottom
    9. CLOSING SLIDE — Professional conclusion
    deck.add_closing_slide(
        speaker_notes="Thank audience, reiterate the ask, invite questions"
    )

    Visual: Tessellated background, centered tagline, URL, logo
    ═══════════════════════════════════════════════════════════════════════════ CONTENTITEM ADVANCED FEATURES ═══════════════════════════════════════════════════════════════════════════════
    ContentItem(
        text="Point text",
        level=1,              # 1=Main, 2=Sub, 3=Detail (indentation + size)
        extra_space=False,    # Add space after this item
        icon="✓",             # Prefix: "✓", "→", "•", "★", "✗"
        highlight=False,      # Apply accent background
        color="#179963"       # Override text color (use theme colors)
    )

    ═══════════════════════════════════════════════════════════════════════════════ AESTHETIC DESIGN GUIDELINES ═══════════════════════════════════════════════════════════════════════════════
    SLIDE TYPE SELECTION MATRIX:
    Content Purpose	Recommended Method	Background
    Opening hook/shocking stat	add_metric_highlight_slide()	Dark
    Problem statement	add_content_slide() or add_comparison_slide()	Dark
    Solution overview	add_content_slide()	Light
    Before/After contrast	add_comparison_slide()	Light
    Market size/opportunity	add_chart_slide("pie")	Light
    Financial projections	add_table_slide() with highlight_row	Light
    Key differentiator	add_section_slide(KEY_MESSAGE)	Dark
    Competitive position	add_two_column_slide()	Light
    The Ask/Investment	add_content_slide()	Light

    COLOR PSYCHOLOGY:
    Dark backgrounds: Urgency, gravity, key messages, problems
    Light backgrounds: Optimism, data clarity, solutions
    Navy (#00325C): Authority, trust
    Red (#C81A28): Problems, urgency, pain points
    Green (#179963): Solutions, growth, success
    Sky Blue (#1B9CD0): Innovation, technology
    VISUAL HIERARCHY:
    Title: 100% visual weight (largest)
    Key metric/headline: 80% weight
    Supporting points: 60% weight
    Details/footnotes: 20% weight
    ═══════════════════════════════════════════════════════════════════════════════ SPEAKER NOTES FORMAT — REQUIRED FOR EVERY SLIDE ═══════════════════════════════════════════════════════════════════════════════
    Every slide MUST include speaker_notes parameter with this structure:

    speaker_notes=\"\"\"
    OPEN: [First sentence to say when slide appears]
    KEY POINT: [The ONE thing audience must remember]
    EVIDENCE: [Data or anecdote to mention verbally]
    TRANSITION: [Bridge to next slide]
    Q&A PREP: [Anticipated question and brief answer]
    \"\"\"

    ═══════════════════════════════════════════════════════════════════════════════ RECOMMENDED SLIDE SEQUENCE (8-12 Slides) ═══════════════════════════════════════════════════════════════════════════════
    #	Slide Purpose	Method	Background
    1	Cover	add_cover_slide()	Dark+AMP
    2	Problem Impact	add_metric_highlight_slide()	Dark
    3	Problem Detail	add_content_slide()	Dark
    4	Solution	add_content_slide()	Light
    5	How It Works	add_comparison_slide()	Light
    6	Market Opportunity	add_chart_slide()	Light
    7	Key Differentiator	add_section_slide(KEY_MESSAGE)	Dark
    8	Traction/Validation	add_table_slide()	Light
    9	The Ask	add_content_slide()	Light
    10	Closing	add_closing_slide()	Dark

    ═══════════════════════════════════════════════════════════════════════════════ PYTHON SCRIPT REQUIREMENTS ═══════════════════════════════════════════════════════════════════════════════
    CRITICAL RULES: • Variable MUST be named: deck • Do NOT include deck.save() • EVERY slide MUST have speaker_notes parameter • Use Confidentiality.{context.get('confidentiality', 'PUBLIC').upper()} for footer
    INITIALIZATION:
    deck = ADIPresentation(
        confidentiality=Confidentiality.PUBLIC,  # or CONFIDENTIAL, INTERNAL_ONLY
        year=2026
    )

    AVAILABLE IN NAMESPACE:
    ADIPresentation, ContentItem, ChartSeries, TableData
    Confidentiality, TableStyle, SectionSlideType
    datetime, date
    CURRENCY: All monetary values in {context['currency']}
    ═══════════════════════════════════════════════════════════════════════════════ OUTPUT REQUIREMENTS ═══════════════════════════════════════════════════════════════════════════════
    PART 1: VERIFICATION REPORT (Brief)
    Inputs enhanced and why
    Statistics needing verification
    Assumptions made
    PART 2: SLIDE-BY-SLIDE CONTENT For each slide: Title, content outline, visual rationale, speaker notes outline
    PART 3: COMPLETE PYTHON SCRIPT Generate full script using ADIPresentation with speaker_notes on EVERY slide.
    IMPORTANT: Do NOT include bibliography or references section.
    Now generate the verification report, slide content, and Python script. """


    def _generate_word_prompt(self, context: dict) -> str:
        """Generate prompt for Deep Dive Summary with section-specific guidance."""
        specific = context.get("specific_inputs", {})
        
        # Build section guidance
        section_guidance = ""
        if any(specific.values()):
            section_guidance = f"""
        
        ═══════════════════════════════════════════════════════════════════════════════
        USER-PROVIDED SECTION CONTENT GUIDANCE
        ═══════════════════════════════════════════════════════════════════════════════

        Use the following guidance to shape each section of the report:

        **EXECUTIVE SUMMARY Guidance:**
        {specific.get('exec_summary', '(No specific guidance provided)')}

        **PURPOSE & SCOPE Guidance:**
        {specific.get('purpose_scope', '(No specific guidance provided)')}

        **TECHNOLOGY & MARKET Guidance:**
        {specific.get('tech_market', '(No specific guidance provided)')}

        **COMPETITIVE LANDSCAPE Guidance:**
        {specific.get('competitive', '(No specific guidance provided)')}

        **RATIONALE FOR GO DECISION Guidance:**
        {specific.get('rationale', '(No specific guidance provided)')}

        **RISKS & UNKNOWNS Guidance:**
        {specific.get('risks', '(No specific guidance provided)')}

        **TECHNOLOGY TRIGGERS & MARKET DYNAMICS Guidance:**
        {specific.get('triggers', '(No specific guidance provided)')}

        **COLLABORATION & WATCH LIST Guidance:**
        {specific.get('collaboration', '(No specific guidance provided)')}

        **NEXT STEPS Guidance:**
        {specific.get('next_steps', '(No specific guidance provided)')}

        IMPORTANT: Incorporate the above guidance into the corresponding sections. Expand on the provided content while maintaining the structured format.
        """
            
        return f"""You are acting as a senior innovation analyst and strategic advisor at Analog Devices' Analog Garage innovation unit, specializing in technology commercialization and investment decisions for the {context['industry']} sector.
                I need you to create a comprehensive Deep Dive Summary document with a recommendation to proceed to an Exploration project for the following innovation domain.
                ═══════════════════════════════════════════════════════════════════════════════ DOCUMENT OBJECTIVE ═══════════════════════════════════════════════════════════════════════════════
                Primary Goal: Produce a structured Deep Dive Summary document that provides a thorough assessment of a technology/market domain and delivers a clear "Recommend to Explore" decision with supporting rationale.
                Target Audience: ADI Innovation leadership, Analog Garage stakeholders, investment committee
                Document Purpose: Support the decision to advance from Deep Dive assessment to Exploration project phase
                ═══════════════════════════════════════════════════════════════════════════════ INNOVATION DOMAIN OVERVIEW ═══════════════════════════════════════════════════════════════════════════════
                Innovation/Domain Name: {context['innovation_name']}
                Target Market/Industry: {context['industry']}
                Geographic Scope: {context['geographic_scope']}
                Innovation Stage: {context['innovation_stage']}
                Analysis Timeframe: {context['analysis_timeframe']}
                ─────────────────────────────────────────────────────────────────────────────── DETAILED DESCRIPTION ───────────────────────────────────────────────────────────────────────────────
                {context['innovation_description']}
                {context['problem_section']}
                {context['customer_section']}
                {context['advantage_section']}
                ─────────────────────────────────────────────────────────────────────────────── MARKET CONTEXT ───────────────────────────────────────────────────────────────────────────────
                {context['market_section']}
                ─────────────────────────────────────────────────────────────────────────────── REGULATORY & IP CONTEXT ───────────────────────────────────────────────────────────────────────────────
                {context['regulatory_section']}
                ─────────────────────────────────────────────────────────────────────────────── KEY CONSIDERATIONS ───────────────────────────────────────────────────────────────────────────────
                {context['risks_section']}
                {context['assumptions_section']}
                ═══════════════════════════════════════════════════════════════════════════════ DEEP DIVE SUMMARY DOCUMENT STRUCTURE ═══════════════════════════════════════════════════════════════════════════════
                Please provide a comprehensive report following this EXACT structure:
                ─────────────────────────────────────────────────────────────────────────────── SECTION: EXECUTIVE SUMMARY ───────────────────────────────────────────────────────────────────────────────
                Brief Overview: One-paragraph summary of the domain studied, the main findings, and the recommendation to proceed to an Exploration project. What is the clearest articulation we have right now of the problem statement and why we believe it's a good fit for ADI?
                Decision Statement: Clear articulation of the "go" decision and its rationale in plain language.
                ─────────────────────────────────────────────────────────────────────────────── SECTION: PURPOSE & SCOPE ───────────────────────────────────────────────────────────────────────────────
                Objective: What was the goal of this assessment? (e.g., evaluate potential for entry, partnership, or investment)
                Exploration Focus: What specific questions or hypotheses will the Exploration phase address?
                Scope:
                What specific subdomains, technologies, or market segments were included/excluded in the study
                What is in-scope for the Exploration phase?
                What is explicitly out-of-scope at this stage? Why?
                Assumptions & Constraints: Any key assumptions (e.g., data availability, access to testbeds), or constraints (budget, timeline, regulatory)
                ─────────────────────────────────────────────────────────────────────────────── SECTION: CURRENT STATE OF TECHNOLOGY & MARKET ───────────────────────────────────────────────────────────────────────────────
                Technology Overview: Concise summary of the current technological landscape (key technologies, maturity, adoption status) - where will exploration focus to drill deeper? PoC scoping?
                Market Overview: Market size, growth trends, major customer segments, and relevant regulatory or economic factors - what are the key questions for exploration, around Proof of Need (PoN)? What problem statements were identified in this study?
                Key Players: Table of major companies, startups, or research groups active in this space and relevance
                ─────────────────────────────────────────────────────────────────────────────── SECTION: COMPETITIVE LANDSCAPE ───────────────────────────────────────────────────────────────────────────────
                Competitor Summary: Brief notes on offerings and market positions, recent moves, flag ones to watch
                Competitive Analysis: High-level SWOT or perceptual map
                Emerging Players: Notable startups or new entrants to watch
                ─────────────────────────────────────────────────────────────────────────────── SECTION: RATIONALE FOR GO DECISION ───────────────────────────────────────────────────────────────────────────────
                Key Assessment Findings: Summarize the most important findings from the Deep Dive
                Summary of Main Reasons to Proceed: (e.g., market opportunity, technology readiness, strategic fit, unique value proposition)
                Supporting Data: Any critical data points, trends, or analysis that support the decision
                ─────────────────────────────────────────────────────────────────────────────── SECTION: RISKS & UNKNOWNS ───────────────────────────────────────────────────────────────────────────────
                Key Risks: Initial identification of major risks and uncertainties
                Mitigation Strategies: Early ideas for how these will be addressed in Exploration
                IP Strategy / Recommendation: Note: Changes to these IP assumptions could invalidate this recommendation
                ─────────────────────────────────────────────────────────────────────────────── SECTION: TECHNOLOGY TRIGGERS ───────────────────────────────────────────────────────────────────────────────
                If these change, they would invalidate or significantly alter the recommendation
                Breakthroughs/Standards: What new technologies or standards would accelerate or change the opportunity?
                ─────────────────────────────────────────────────────────────────────────────── SECTION: MARKET DYNAMICS ───────────────────────────────────────────────────────────────────────────────
                Customer/Regulatory Shifts: What changes in demand, regulation, or competition would impact the project?
                ─────────────────────────────────────────────────────────────────────────────── SECTION: COLLABORATION OPPORTUNITIES ───────────────────────────────────────────────────────────────────────────────
                Potential Partners/Consortia: Are there ecosystem changes or new partnerships that could enhance the project?
                ─────────────────────────────────────────────────────────────────────────────── SECTION: COMPANIES AND TRENDS TO WATCH ───────────────────────────────────────────────────────────────────────────────
                Key Companies: Incumbents & Startups: List of companies to monitor, with a brief note on relevance
                Trends/Signals: Notable Trends/Events: Funding rounds, regulatory changes, technology launches, etc., that should trigger a review or pivot
                ─────────────────────────────────────────────────────────────────────────────── SECTION: FOLLOW-UP & KNOWLEDGE SHARING ───────────────────────────────────────────────────────────────────────────────
                Next Steps: Immediate Actions: What are the next steps to launch the Exploration project?
                Knowledge Sharing: Report Storage: Where will this report and future updates be stored?
                Contact: Project Lead/Point of Contact: Who to reach for questions or further discussion
                ═══════════════════════════════════════════════════════════════════════════════ OUTPUT FORMAT ═══════════════════════════════════════════════════════════════════════════════
                All monetary values in {context['currency']}.
                Document Formatting Requirements:
                Use clear section headers matching the structure above
                Tables should be properly formatted
                Bullet points for lists
                Bold for emphasis on key terms
                Professional, executive-ready tone
                IMPORTANT: Do NOT include a bibliography, references, or sources section at the end of your response.
                ═══════════════════════════════════════════════════════════════════════════════ PYTHON-DOCX SCRIPT GENERATION ═══════════════════════════════════════════════════════════════════════════════
                After providing the report content above, generate a Python script using python-docx to create the Word document.
                CRITICAL RULES: • Document variable MUST be named exactly: doc • Do NOT include doc.save() — the application handles file saving • Use only python-docx library • Match the formatting of the "Deep Dive Summary - Recommend to Explore" template
                Document Styling Requirements: • Title: "DeepDive Summary" — Large, bold, dark blue (RGB: 0, 103, 185) • Section Headers: Bold, dark blue, larger font (14-16pt) • Subsection Headers: Bold, slightly smaller (12pt) • Body Text: Regular, 11pt • Tables: Light blue header row, alternating row colors • Professional spacing between sections
                Available Imports:

                from docx import Document
                from docx.shared import Inches, Pt, RGBColor, Cm
                from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
                from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ROW_HEIGHT_RULE
                from docx.enum.style import WD_STYLE_TYPE
                from docx.oxml.ns import qn, nsdecls
                from docx.oxml import parse_xml, OxmlElement

            """

    def _show_generation_confirmation(self, output_config: dict, values: dict):
        """Show output-specific confirmation message."""
        output_type = output_config.get("id", "excel_value_model")
        icon = output_config.get("icon", "📊")
        name = output_config.get("name", "Output")
        
        if output_type == "excel_value_model":
            message = (
                f"{icon} {name} prompt generated!\n\n"
                f"• 100% penetration analysis\n"
                f"• Geographic breakdown included\n"
                f"• 2026-2040 projections with growth rates\n\n"
                f"Click 'Copy to Clipboard' to use with your AI tool."
            )
        elif output_type == "powerpoint_pitch":
            message = (
                f"{icon} {name} prompt generated!\n\n"
                f"• 12-slide executive pitch structure\n"
                f"• Speaker notes included\n"
                f"• Visual suggestions for each slide\n\n"
                f"Click 'Copy to Clipboard' to use with your AI tool."
            )
        elif output_type == "word_gonogo":
            message = (
                f"{icon} {name} prompt generated!\n\n"
                f"• Comprehensive Deep Dive structure\n"
                f"• Aligned with Analog Garage template\n"
                f"• Exploration recommendation format\n\n"
                f"Click 'Copy to Clipboard' to use with your AI tool."
            )
        else:
            message = f"{icon} {name} prompt generated!\n\nClick 'Copy to Clipboard' to use with your AI tool."
        
        messagebox.showinfo("Prompt Generated", message)

    def _copy_prompt(self):
        """Copy prompt to clipboard."""
        content = self.prompt_output.get_content()
        if not content:
            messagebox.showwarning("Warning", "No prompt to copy. Generate a prompt first.")
            return
        
        self.root.clipboard_clear()
        self.root.clipboard_append(content)
        self.session.info("Prompt copied to clipboard")
        self._update_status("✓ Prompt copied to clipboard!", ADI_COLORS["success_green"])
        messagebox.showinfo("Copied", "Prompt copied to clipboard!\n\nPaste it into Claude, ChatGPT, or your preferred AI tool.")

    def _export_prompt(self):
        """Export prompt to file."""
        content = self.prompt_output.get_content()
        if not content:
            messagebox.showwarning("Warning", "No prompt to export.")
            return
        
        context = context_service.get_current() or {}
        output_config = self.output_selector.get_output_type_config()
        output_type = output_config.get("short_name", "Prompt").replace(" ", "")
        
        name = context.get("innovation_name", "prompt")
        safe_name = "".join(c for c in name if c.isalnum() or c in (' ', '-', '_')).rstrip()[:30]
        safe_name = safe_name.replace(' ', '_')
        default_filename = f"{output_type}_{safe_name}_{datetime.now().strftime('%Y%m%d')}.txt"
        
        file_path = filedialog.asksaveasfilename(
            defaultextension=".txt",
            filetypes=[("Text files", "*.txt"), ("All files", "*.*")],
            title="Export Prompt",
            initialfile=default_filename
        )
        
        if file_path:
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content)
            self.session.info(f"Prompt exported: {file_path}")
            self._update_status(f"✓ Exported: {os.path.basename(file_path)}", ADI_COLORS["success_green"])

    def _clear_form(self):
        """Clear the input form."""
        self.form.clear()
        self.prompt_output.clear()
        context_service.clear_current()
        self._update_status("Form cleared", ADI_COLORS["text_gray"])

    # =========================================================================
    # SCRIPT EXECUTION - Updated for multiple output types
    # =========================================================================

    def _execute_script(self):
        """Execute script based on current output type."""
        script = self.script_input.get_script()
        
        if not script:
            messagebox.showwarning("Warning", "Please paste a Python script first.")
            return
        
        output_config = self.output_selector.get_output_type_config()
        script_type = output_config.get("script_type", "excel")
        
        self._update_status("Executing script...", ADI_COLORS["warning_yellow"])
        self.session.set_script(script)
        
        # Execute based on script type
        if script_type == "excel":
            self._execute_excel_script(script)
        elif script_type == "powerpoint":
            self._execute_powerpoint_script(script)
        elif script_type == "word":
            self._execute_word_script(script)
        else:
            messagebox.showerror("Error", f"Unknown script type: {script_type}")

    def _execute_excel_script(self, script: str):
        """Execute Excel (openpyxl) script."""
        result = excel_executor.execute(script)
        
        if result.success and result.workbook:
            self.session.success("Excel script executed successfully")
            
            file_path = filedialog.asksaveasfilename(
                defaultextension=".xlsx",
                filetypes=[("Excel files", "*.xlsx"), ("All files", "*.*")],
                title="Save Excel Workbook"
            )
            
            if file_path:
                success, error = excel_executor.save_workbook(result.workbook, file_path)
                if success:
                    self.session.log_file_saved(file_path)
                    self._update_status(f"✓ Saved: {os.path.basename(file_path)}", ADI_COLORS["success_green"])
                    messagebox.showinfo("Success", f"Excel workbook saved to:\n{file_path}")
                else:
                    self.session.error(error)
                    messagebox.showerror("Save Error", error)
            else:
                self._update_status("Save cancelled", ADI_COLORS["text_gray"])
        else:
            error_details = result.traceback if result.traceback else result.error_message
            self.session.set_error(error_details)
            self._update_status("Execution error", ADI_COLORS["error_red"])
            ErrorDialogComponent(self.root, error_details, on_export_log=self._export_log)

    def _execute_powerpoint_script(self, script: str):
        """Execute PowerPoint (python-pptx / ADIPresentation) script."""
        try:
            # Use the powerpoint_processor module which handles 'deck' variable
            result = execute_pptx_script(script)
            
            if result.has_error:
                error_details = result.traceback if result.traceback else result.error_message
                self.session.set_error(error_details)
                self._update_status("Execution error", ADI_COLORS["error_red"])
                ErrorDialogComponent(self.root, error_details, on_export_log=self._export_log)
                return
            
            # Get the presentation object
            presentation = result.presentation
            
            # Prompt user for save location
            file_path = filedialog.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")],
                title="Save PowerPoint Presentation"
            )
            
            if file_path:
                success, error = save_pptx(presentation, file_path)
                if success:
                    self.session.log_file_saved(file_path)
                    self._update_status(f"✓ Saved: {os.path.basename(file_path)}", ADI_COLORS["success_green"])
                    messagebox.showinfo("Success", f"PowerPoint presentation saved to:\n{file_path}")
                else:
                    self.session.log_pptx_error(error)
                    messagebox.showerror("Save Error", f"Failed to save: {error}")
            else:
                self._update_status("Save cancelled", ADI_COLORS["text_gray"])
                
        except ImportError as e:
            messagebox.showerror(
                "Missing Library",
                f"python-pptx is required for PowerPoint generation.\n\nInstall with: pip install python-pptx\n\nError: {e}"
            )
        except Exception as e:
            error_details = traceback.format_exc()
            self.session.set_error(error_details)
            self._update_status("Execution error", ADI_COLORS["error_red"])
            ErrorDialogComponent(self.root, error_details, on_export_log=self._export_log)


    def _execute_word_script(self, script: str):
        """Execute Word (python-docx) script."""
        try:
            # Validate script
            if "doc" not in script and "Document" not in script:
                raise ValueError("Script must create a document object named 'doc'")
            
            if "doc.save(" in script:
                raise ValueError("doc.save() should not be included - the app handles saving")
            
            # Build execution namespace for python-docx
            namespace = self._build_docx_namespace()
            
            exec(script, namespace, namespace)
            
            if 'doc' in namespace:
                doc = namespace['doc']
                
                file_path = filedialog.asksaveasfilename(
                    defaultextension=".docx",
                    filetypes=[("Word files", "*.docx"), ("All files", "*.*")],
                    title="Save Word Document"
                )
                
                if file_path:
                    doc.save(file_path)
                    self.session.log_file_saved(file_path)
                    self._update_status(f"✓ Saved: {os.path.basename(file_path)}", ADI_COLORS["success_green"])
                    messagebox.showinfo("Success", f"Word document saved to:\n{file_path}")
                else:
                    self._update_status("Save cancelled", ADI_COLORS["text_gray"])
            else:
                raise ValueError("Script did not create a document object named 'doc'")
                
        except ImportError as e:
            messagebox.showerror(
                "Missing Library",
                f"python-docx is required for Word generation.\n\nInstall with: pip install python-docx\n\nError: {e}"
            )
        except Exception as e:
            error_details = traceback.format_exc()
            self.session.set_error(error_details)
            self._update_status("Execution error", ADI_COLORS["error_red"])
            ErrorDialogComponent(self.root, error_details, on_export_log=self._export_log)

    def _build_pptx_namespace(self) -> dict:
        """Build namespace for python-pptx execution."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE
            
            namespace = {
                'Presentation': Presentation,
                'Inches': Inches,
                'Pt': Pt,
                'RGBColor': RGBColor,
                'PP_ALIGN': PP_ALIGN,
                'MSO_ANCHOR': MSO_ANCHOR,
                'MSO_SHAPE': MSO_SHAPE,
            }
            
            # Try to add chart support
            try:
                from pptx.enum.chart import XL_CHART_TYPE
                from pptx.chart.data import CategoryChartData
                namespace['XL_CHART_TYPE'] = XL_CHART_TYPE
                namespace['CategoryChartData'] = CategoryChartData
            except ImportError:
                pass
            
            return namespace
        except ImportError:
            raise ImportError("python-pptx is not installed")

    def _build_docx_namespace(self) -> dict:
        """Build namespace for python-docx execution."""
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor, Cm
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.table import WD_TABLE_ALIGNMENT
            from docx.enum.style import WD_STYLE_TYPE
            
            namespace = {
                'Document': Document,
                'Inches': Inches,
                'Pt': Pt,
                'Cm': Cm,
                'RGBColor': RGBColor,
                'WD_ALIGN_PARAGRAPH': WD_ALIGN_PARAGRAPH,
                'WD_TABLE_ALIGNMENT': WD_TABLE_ALIGNMENT,
                'WD_STYLE_TYPE': WD_STYLE_TYPE,
            }
            
            # Try to add additional imports
            try:
                from docx.enum.text import WD_LINE_SPACING
                from docx.enum.table import WD_ROW_HEIGHT_RULE
                from docx.oxml.ns import qn, nsdecls
                from docx.oxml import parse_xml, OxmlElement
                namespace['WD_LINE_SPACING'] = WD_LINE_SPACING
                namespace['WD_ROW_HEIGHT_RULE'] = WD_ROW_HEIGHT_RULE
                namespace['qn'] = qn
                namespace['nsdecls'] = nsdecls
                namespace['parse_xml'] = parse_xml
                namespace['OxmlElement'] = OxmlElement
            except ImportError:
                pass
            
            return namespace
        except ImportError:
            raise ImportError("python-docx is not installed")

    def _clear_script(self):
        """Clear script input."""
        self.script_input.clear()
        self._update_status("Script cleared", ADI_COLORS["text_gray"])

    # =========================================================================
    # OUTPUT TYPE SELECTION
    # =========================================================================

    def _select_output_type(self, output_id: str):
        """Handle output type selection from menu."""
        self.output_selector.set_output_type(output_id)

    def _select_template(self, template_id: str):
        """Handle template selection from menu."""
        template = template_registry.get(template_id)
        if template:
            prompt_service.default_template_id = template_id
            self.session.info(f"Template selected: {template.name}")
            self._update_status(f"Template: {template.name}", ADI_COLORS["primary_blue"])
            messagebox.showinfo("Template Selected", f"Now using: {template.name}\n\n{template.description}")

    # =========================================================================
    # UTILITY METHODS
    # =========================================================================

    def _update_status(self, message: str, color: str = None):
        """Update status bar."""
        self.status_bar.update_status(message, color)

    def _export_log(self):
        """Export session log."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        default_filename = f"AnalogGarage_Log_{timestamp}.txt"
        
        file_path = filedialog.asksaveasfilename(
            defaultextension=".txt",
            filetypes=[("Text files", "*.txt"), ("Log files", "*.log")],
            title="Export Session Log",
            initialfile=default_filename
        )
        
        if file_path:
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(self.session.generate_report())
            self._update_status(f"Log exported: {os.path.basename(file_path)}", ADI_COLORS["success_green"])
            messagebox.showinfo("Exported", f"Session log saved to:\n{file_path}")

    def _show_help(self):
        """Show help dialog."""
        help_text = """
            ANALOG GARAGE WORKBENCH - User Guide (v3.0) ════════════════════════════════════════════════
            This tool generates prompts for three types of deliverables:
            OUTPUT TYPES: ───────────── 📊 Excel Value Creation Model • Comprehensive financial model • 2026-2040 projections with growth rates • Charts and sensitivity analysis
            📽️ Executive Pitch Deck • 12-slide PowerPoint structure • Speaker notes included • Visual suggestions
            📋 Deep Dive Summary - Recommend to Explore • Comprehensive analysis document • Aligned with Analog Garage template • Exploration recommendation format
            TAB 1: GENERATE PROMPT ─────────────────────
            Select output type from dropdown
            Fill in innovation details
            Click "Generate Prompt"
            Copy and paste into AI tool (Claude, ChatGPT, etc.)
            TAB 2: EXECUTE SCRIPT ─────────────────────
            Paste AI-generated script
            Click Execute button
            Save your file
            SCRIPT REQUIREMENTS BY TYPE: ─────────────────────────── • Excel: Use 'wb' for workbook (openpyxl)
            Do NOT include wb.save()
            • PowerPoint: Use 'prs' for presentation (python-pptx)
            Do NOT include prs.save()
            • Word: Use 'doc' for document (python-docx)
            Do NOT include doc.save()
            KEYBOARD SHORTCUTS: ────────────────── • Cmd/Ctrl + Q: Quit
            REQUIRED LIBRARIES: ────────────────── • openpyxl (for Excel) • python-pptx (for PowerPoint) - pip install python-pptx • python-docx (for Word) - pip install python-docx 
            """
        self._show_help_window(help_text)
                
    def _show_help_window(self, help_text: str):
        """Display help window."""
        help_window = tk.Toplevel(self.root)
        help_window.title("User Guide")
        help_window.geometry("700x650")
        help_window.configure(bg=ADI_COLORS["white"])
            
        text = scrolledtext.ScrolledText(
            help_window,
            font=("Arial", 11),
            bg=ADI_COLORS["white"],
            wrap=tk.WORD,
            padx=20,
            pady=15
        )
        text.insert(tk.END, help_text)
        text.config(state=tk.DISABLED)
        text.pack(fill=tk.BOTH, expand=True)
        
        tk.Button(
            help_window,
            text="Close",
            command=help_window.destroy,
            font=("Arial", 10),
            padx=20,
            pady=5
        ).pack(pady=15)

    def _show_about(self):
        """Show about dialog."""
        messagebox.showinfo(
            "About",
            f"{APP_NAME}\n"
            f"Version 3.0\n\n"
            "Multi-Output Innovation Analysis Tool:\n\n"
            "📊 Excel Value Creation Models\n"
            "📽️ Executive Pitch Decks\n"
            "📋 Deep Dive Summary Reports\n\n"
            "Features:\n"
            "• Multiple output type support\n"
            "• Template-based prompt generation\n"
            "• Script execution for Office files\n"
            "• Session logging\n\n"
            "© 2026 Analog Devices, Inc.\n"
            "Analog Garage - HC & OIT Innovation"
        )

    def generate_powerpoint(self, script: str, output_path: str) -> Tuple[bool, str]:
        """
        Generate PowerPoint from script.
        
        Args:
            script: Python script that creates 'deck' variable
            output_path: Where to save the .pptx file
            
        Returns:
            Tuple of (success, message)
        """
        self.session.set_pptx_script(script)
        self.session.info("Executing PowerPoint generation script...")
        
        # Execute script
        result = execute_pptx_script(script)
        
        if result.has_error:
            self.session.log_pptx_error(result.error_message)
            return False, f"Script error: {result.error_message}"
        
        # Save presentation
        success, error = save_pptx(result.presentation, output_path)
        
        if success:
            slide_count = result.presentation.slide_count
            self.session.log_pptx_generated(slide_count, output_path)
            return True, f"Successfully created {slide_count}-slide presentation: {output_path}"
        else:
            self.session.log_pptx_error(error)
            return False, f"Save error: {error}"
    
    def generate_powerpoint_from_data(self, 
                                       title: str,
                                       slides_data: List[Dict],
                                       output_path: str,
                                       confidentiality: str = "public") -> Tuple[bool, str]:
        """
        Generate PowerPoint from structured data (no script needed).
        
        Args:
            title: Presentation title
            slides_data: List of slide definitions
            output_path: Where to save the .pptx file
            confidentiality: 'public', 'confidential', or 'internal_only'
            
        Returns:
            Tuple of (success, message)
        """
        try:
            self.session.info(f"Generating PowerPoint: {title}")
            
            deck = quick_presentation(
                title=title,
                confidentiality=confidentiality,
                slides=slides_data
            )
            
            filepath = deck.save(output_path)
            
            self.session.log_pptx_generated(deck.slide_count, filepath)
            return True, f"Successfully created: {filepath}"
            
        except Exception as e:
            self.session.log_pptx_error(str(e))
            return False, f"Error: {str(e)}"
    
    def get_pptx_help(self) -> str:
        """Get help text for PowerPoint script writing."""
        namespace_info = get_pptx_namespace_info()
        
        lines = [
            "=" * 60,
            "POWERPOINT SCRIPT HELP",
            "=" * 60,
            "",
            "Available Classes:",
            *[f"  - {c}" for c in sorted(namespace_info["classes"])],
            "",
            "Available Functions:",
            *[f"  - {f}" for f in sorted(namespace_info["functions"])],
            "",
            "Available Enums:",
            *[f"  - {e}" for e in sorted(namespace_info["enums"])],
            "",
            "Example Script:",
            "  deck = ADIPresentation(confidentiality=Confidentiality.CONFIDENTIAL)",
            "  deck.add_cover_slide('My Presentation', 'Subtitle')",
            "  deck.add_content_slide('Agenda', [ContentItem('Item 1', 1)])",
            "  deck.add_closing_slide()",
            "",
            "NOTE: Do not include save() - the application handles saving.",
            "=" * 60,
        ]
        
        return "\n".join(lines)

    def run(self):
        """Start the application."""
        self.session.info("Application started")
        self.root.mainloop()

# =============================================================================
# ENTRY POINT
# =============================================================================

def main():
    """Application entry point."""
    app = AnalogGarageWorkbench()
    app.run()

if __name__ == "__main__":
    main()
