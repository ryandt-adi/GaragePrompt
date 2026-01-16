#!/usr/bin/env python3
"""
ADI PowerPoint Generator v4.0 — Full Template Compliance
=========================================================
FIXED:
- RgbColor import (was RGBColor)
- Added proper error handling
- Fixed type hints
- Added missing imports

Version: 4.0.1
"""

from pptx import Presentation
from pptx.util import Pt, Cm, Emu
# FIXED: Correct import is RgbColor, not RGBColor
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pathlib import Path
from datetime import datetime
from typing import List, Optional, Union, Tuple, Dict, Any
from dataclasses import dataclass
import warnings
import os

# Import centralized configuration
try:
    from adi_template_config import (
        COLORS, DIMENSIONS, TYPOGRAPHY, CONTAINERS, AMP_CONFIG,
        FOOTER_TEMPLATES, TITLE_CASE_LOWERCASE_WORDS, CLOSING_TAGLINE,
        COMPANY_URL, TABLE_CONFIG, ASSET_PATHS, VALIDATION,
        Confidentiality, SlideType, TableStyle, SectionSlideType,
        ChartColorScheme, ADIColorPalette
    )
except ImportError as e:
    raise ImportError(f"Could not import adi_template_config: {e}. Ensure the module exists.")

# Optional PIL import
try:
    from PIL import Image, ImageDraw
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    warnings.warn("PIL not installed. AMP overlay generation disabled. Install with: pip install Pillow")


# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class ContentItem:
    """Content item with text hierarchy level."""
    text: str
    level: int = 1
    extra_space: bool = False
    icon: Optional[str] = None  # NEW: "✓", "→", "•", "★", etc.
    highlight: bool = False      # NEW: Apply accent background
    color: Optional[str] = None  # NEW: Override text color


    def __post_init__(self):
        if self.level not in [1, 2, 3]:
            raise ValueError(f"Level must be 1, 2, or 3. Got {self.level}")


@dataclass
class ChartSeries:
    """Chart data series."""
    name: str
    values: List[float]


@dataclass
class TableData:
    """Table with headers and rows."""
    headers: List[str]
    rows: List[List[str]]
    highlight_row: Optional[int] = None  # NEW: Index of row to highlight
    value_formatting: Optional[Dict] = None  # NEW: {"col_index": "currency"|"percent"|"trend"}


    def __post_init__(self):
        for i, row in enumerate(self.rows):
            if len(row) != len(self.headers):
                raise ValueError(f"Row {i} has {len(row)} columns but headers have {len(self.headers)}")


@dataclass
class ValidationResult:
    """Validation check result."""
    passed: bool
    message: str
    severity: str = "warning"


# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """
    Convert hex color to RgbColor.
    FIXED: Returns RgbColor (correct class name)
    """
    h = hex_color.lstrip('#')
    if len(h) != 6:
        raise ValueError(f"Invalid hex color: {hex_color}")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def to_title_case(text: str) -> str:
    """Apply ADI title case rules."""
    if not text:
        return text
    
    words = text.split()
    result = []
    
    for i, word in enumerate(words):
        if word.isupper() and len(word) > 1:
            result.append(word)
        elif i == 0:
            result.append(word.capitalize())
        elif word.lower() in TITLE_CASE_LOWERCASE_WORDS and len(word) <= 4:
            result.append(word.lower())
        else:
            result.append(word.capitalize())
    
    return ' '.join(result)


def format_bullet(text: str, level: int) -> str:
    """Format text with bullet character for levels 2-3."""
    if level == 1 or not text:
        return text
    return f"{TYPOGRAPHY.BULLET_CHAR} {text}"


# =============================================================================
# ASSET MANAGER
# =============================================================================

class AssetManager:
    """Manages brand asset files."""
    
    def __init__(self, base_path: str = "assets"):
        self.base = Path(base_path)
        # Create directory if it doesn't exist
        try:
            self.base.mkdir(parents=True, exist_ok=True)
        except Exception as e:
            warnings.warn(f"Could not create assets directory: {e}")
    
    @property
    def logo_blue(self) -> Path:
        return self.base / ASSET_PATHS.LOGO_BLUE
    
    @property
    def logo_white(self) -> Path:
        return self.base / ASSET_PATHS.LOGO_WHITE
    
    @property
    def cover_background(self) -> Path:
        return self.base / ASSET_PATHS.COVER_BACKGROUND
    
    @property
    def tessellated_background(self) -> Path:
        return self.base / ASSET_PATHS.TESSELLATED_BACKGROUND
    
    @property
    def amp_overlay(self) -> Path:
        return self.base / ASSET_PATHS.AMP_OVERLAY
    
    def get_missing(self) -> List[str]:
        """Return list of missing assets."""
        missing = []
        for attr in ['logo_blue', 'logo_white', 'cover_background', 'tessellated_background']:
            path = getattr(self, attr)
            if not path.exists():
                missing.append(str(path))
        return missing
    
    def ensure_amp_overlay(self) -> bool:
        """Generate AMP overlay if missing."""
        if self.amp_overlay.exists():
            return True
        
        if not HAS_PIL:
            return False
        
        return AMPOverlayGenerator.create(self.amp_overlay)


# =============================================================================
# AMP OVERLAY GENERATOR
# =============================================================================

class AMPOverlayGenerator:
    """Generates the AMP triangle overlay as transparent PNG."""
    
    @staticmethod
    def create(output_path: Path, width: int = 3387, height: int = 1905) -> bool:
        """Create transparent PNG with exact AMP triangle shape."""
        if not HAS_PIL:
            warnings.warn("PIL required for AMP overlay generation")
            return False
        
        try:
            img = Image.new('RGBA', (width, height), (0, 0, 0, 0))
            draw = ImageDraw.Draw(img)
            
            vertices = AMP_CONFIG.get_vertices_pixels(width, height)
            
            color_hex = AMP_CONFIG.COLOR.lstrip('#')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            
            draw.polygon(vertices, fill=(r, g, b, 255))
            
            # Ensure parent directory exists
            output_path.parent.mkdir(parents=True, exist_ok=True)
            img.save(str(output_path), 'PNG')
            
            return True
        except Exception as e:
            warnings.warn(f"Failed to create AMP overlay: {e}")
            return False


# =============================================================================
# BRAND VALIDATOR
# =============================================================================

class BrandValidator:
    """Validates content against ADI brand guidelines."""
    
    def __init__(self):
        self.results: List[ValidationResult] = []
    
    def clear(self):
        """Clear all validation results."""
        self.results = []
    
    def add(self, passed: bool, message: str, severity: str = "warning"):
        """Add a validation result."""
        self.results.append(ValidationResult(passed, message, severity))
    
    def validate_title_size(self, size_pt: int) -> bool:
        """Validate title font size."""
        if size_pt < VALIDATION.MIN_TITLE_SIZE_PT:
            self.add(False, f"Title size {size_pt}pt below minimum {VALIDATION.MIN_TITLE_SIZE_PT}pt", "error")
            return False
        return True
    
    def validate_slide_count(self, count: int) -> bool:
        """Validate slide count against recommendation."""
        if count > VALIDATION.MAX_RECOMMENDED_SLIDES:
            self.add(False, f"Slide count ({count}) exceeds recommended {VALIDATION.MAX_RECOMMENDED_SLIDES}", "warning")
            return False
        self.add(True, f"Slide count OK ({count})")
        return True
    
    def validate_table_row_height(self, height_cm: float) -> bool:
        """Validate table row height."""
        if height_cm < VALIDATION.MIN_TABLE_ROW_HEIGHT_CM:
            self.add(False, f"Table row height {height_cm}cm below minimum {VALIDATION.MIN_TABLE_ROW_HEIGHT_CM}cm", "warning")
            return False
        return True
    
    def has_errors(self) -> bool:
        """Check if any errors were found."""
        return any(not r.passed and r.severity == "error" for r in self.results)
    
    def has_warnings(self) -> bool:
        """Check if any warnings were found."""
        return any(not r.passed and r.severity == "warning" for r in self.results)
    
    def get_report(self) -> str:
        """Generate validation report."""
        lines = ["=" * 60, "ADI BRAND COMPLIANCE REPORT", "=" * 60, ""]
        
        errors = [r for r in self.results if not r.passed and r.severity == "error"]
        warns = [r for r in self.results if not r.passed and r.severity == "warning"]
        passed = [r for r in self.results if r.passed]
        
        if errors:
            lines.append("❌ ERRORS:")
            for r in errors:
                lines.append(f"   • {r.message}")
            lines.append("")
        
        if warns:
            lines.append("⚠️  WARNINGS:")
            for r in warns:
                lines.append(f"   • {r.message}")
            lines.append("")
        
        lines.append(f"✓ Passed: {len(passed)} | ⚠ Warnings: {len(warns)} | ❌ Errors: {len(errors)}")
        lines.append("=" * 60)
        
        return "\n".join(lines)


# =============================================================================
# MAIN PRESENTATION CLASS
# =============================================================================

class ADIPresentation:
    """Generate ADI-compliant PowerPoint presentations."""
    
    def __init__(self,
                 confidentiality: Confidentiality = Confidentiality.PUBLIC,
                 year: Optional[int] = None,
                 assets_path: str = "assets"):
        """Initialize presentation."""
        self.prs = Presentation()
        self.confidentiality = confidentiality
        self.year = year or datetime.now().year
        self.assets = AssetManager(assets_path)
        self.validator = BrandValidator()
        self.slide_count = 0
        self.warnings: List[str] = []
        
        # Set slide dimensions
        self.prs.slide_width = Cm(DIMENSIONS.WIDTH_CM)
        self.prs.slide_height = Cm(DIMENSIONS.HEIGHT_CM)
        
        # Generate AMP overlay
        if not self.assets.ensure_amp_overlay():
            self.warnings.append("Could not generate AMP overlay")
        
        # Check for missing assets
        missing = self.assets.get_missing()
        if missing:
            self.warnings.append(f"Missing assets: {', '.join(missing)}")
    
    def _blank_layout(self):
        """Get blank slide layout."""
        return self.prs.slide_layouts[6]
    
    def _get_footer_text(self) -> str:
        """Get footer text for current confidentiality."""
        return FOOTER_TEMPLATES[self.confidentiality].format(year=self.year)
    
    def _get_container(self, key: str) -> Dict[str, float]:
        """Get locked container position."""
        container = CONTAINERS.get(key)
        if not container:
            warnings.warn(f"Unknown container key: {key}")
            return {"left": 0, "top": 0, "width": 10, "height": 5}
        return container
    
    # -------------------------------------------------------------------------
    # BACKGROUND METHODS
    # -------------------------------------------------------------------------
    
    def _add_solid_background(self, slide, color_hex: str):
        """Add solid color background."""
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(color_hex)
        except Exception as e:
            self.warnings.append(f"Failed to set background color: {e}")
    
    def _add_image_background(self, slide, image_path: Path) -> bool:
        """Add full-bleed background image."""
        if not image_path.exists():
            return False
        
        try:
            pic = slide.shapes.add_picture(
                str(image_path),
                left=0, top=0,
                width=Cm(DIMENSIONS.WIDTH_CM),
                height=Cm(DIMENSIONS.HEIGHT_CM)
            )
            
            # Send to back
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(2, pic._element)
            return True
        except Exception as e:
            self.warnings.append(f"Failed to add background image: {e}")
            return False
    
    def _add_amp_overlay(self, slide) -> bool:
        """Add AMP triangle overlay."""
        if self.assets.amp_overlay.exists():
            try:
                slide.shapes.add_picture(
                    str(self.assets.amp_overlay),
                    left=0, top=0,
                    width=Cm(DIMENSIONS.WIDTH_CM),
                    height=Cm(DIMENSIONS.HEIGHT_CM)
                )
                return True
            except Exception as e:
                self.warnings.append(f"Failed to add AMP overlay: {e}")
        
        return self._add_amp_shapes(slide)
    
    def _add_amp_shapes(self, slide) -> bool:
        """Create AMP using shapes (fallback)."""
        try:
            # Rectangle
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Cm(0), Cm(0),
                Cm(AMP_CONFIG.RECT_WIDTH_CM),
                Cm(DIMENSIONS.HEIGHT_CM)
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = hex_to_rgb(AMP_CONFIG.COLOR)
            rect.line.fill.background()
            
            # Triangle
            tri = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                Cm(AMP_CONFIG.RECT_WIDTH_CM), Cm(0),
                Cm(AMP_CONFIG.TRIANGLE_WIDTH_CM),
                Cm(DIMENSIONS.HEIGHT_CM)
            )
            tri.fill.solid()
            tri.fill.fore_color.rgb = hex_to_rgb(AMP_CONFIG.COLOR)
            tri.line.fill.background()
            
            return True
        except Exception as e:
            self.warnings.append(f"Failed to create AMP shapes: {e}")
            return False
    
    # -------------------------------------------------------------------------
    # ELEMENT METHODS
    # -------------------------------------------------------------------------
    
    def _add_logo(self, slide, position: str, white: bool = False):
        """Add logo at specified position."""
        logo_path = self.assets.logo_white if white else self.assets.logo_blue
        if not logo_path.exists():
            return None
        
        pos = self._get_container(f"logo_{position}")
        if not pos or "width" not in pos:
            return None
        
        try:
            return slide.shapes.add_picture(
                str(logo_path),
                Cm(pos.get("left", 0)),
                Cm(pos.get("top", 0)),
                width=Cm(pos.get("width", 4))
            )
        except Exception as e:
            self.warnings.append(f"Failed to add logo: {e}")
            return None
    
    def _add_title(self, slide, title: str, container_key: str,
                   dark_bg: bool = False, size_pt: int = None):
        """Add title at locked container position."""
        size_pt = size_pt or TYPOGRAPHY.TITLE_MIN_PT
        self.validator.validate_title_size(size_pt)
        
        pos = self._get_container(container_key)
        
        try:
            box = slide.shapes.add_textbox(
                Cm(pos["left"]), Cm(pos["top"]),
                Cm(pos["width"]), Cm(pos["height"])
            )
            
            tf = box.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = to_title_case(title)
            p.font.name = TYPOGRAPHY.TITLE_FONT
            p.font.size = Pt(size_pt)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(
                COLORS.TEXT_LIGHT if dark_bg else COLORS.TEXT_DARK
            )
            
            return box
        except Exception as e:
            self.warnings.append(f"Failed to add title: {e}")
            return None
    
    def _add_body_content(self, slide, content: List[ContentItem],
                          container_key: str, dark_bg: bool = False):
        """Add body content with text hierarchy."""
        pos = self._get_container(container_key)
        
        try:
            box = slide.shapes.add_textbox(
                Cm(pos["left"]), Cm(pos["top"]),
                Cm(pos["width"]), Cm(pos["height"])
            )
            
            tf = box.text_frame
            tf.word_wrap = True
            tf.clear()
            
            text_color = COLORS.TEXT_LIGHT if dark_bg else COLORS.TEXT_DARK
            
            level_config = {
                1: (TYPOGRAPHY.LEVEL_1_SIZE, TYPOGRAPHY.LEVEL_1_SPACE_AFTER),
                2: (TYPOGRAPHY.LEVEL_2_SIZE, TYPOGRAPHY.LEVEL_2_SPACE_AFTER),
                3: (TYPOGRAPHY.LEVEL_3_SIZE, TYPOGRAPHY.LEVEL_3_SPACE_AFTER),
            }
            
            for i, item in enumerate(content):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                
                size, space = level_config.get(item.level, level_config[1])
                
                p.text = format_bullet(item.text, item.level)
                p.font.name = TYPOGRAPHY.BODY_FONT
                p.font.size = Pt(size)
                p.font.color.rgb = hex_to_rgb(text_color)
                p.space_after = Pt(space * (2 if item.extra_space else 1))
                
                if item.level > 1:
                    p.level = item.level - 1
            
            return box
        except Exception as e:
            self.warnings.append(f"Failed to add body content: {e}")
            return None
    
    def _add_url(self, slide, container_key: str, dark_bg: bool = True):
        """Add company URL."""
        pos = self._get_container(container_key)
        
        try:
            box = slide.shapes.add_textbox(
                Cm(pos["left"]), Cm(pos["top"]),
                Cm(pos["width"]), Cm(pos["height"])
            )
            
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = COMPANY_URL
            p.font.name = TYPOGRAPHY.BODY_FONT
            p.font.size = Pt(TYPOGRAPHY.URL_SIZE)
            p.font.color.rgb = hex_to_rgb(
                COLORS.TEXT_LIGHT if dark_bg else COLORS.TEXT_DARK
            )
            
            return box
        except Exception as e:
            self.warnings.append(f"Failed to add URL: {e}")
            return None
    
    def _add_footer(self, slide, include_number: bool = True, dark_bg: bool = False):
        """Add footer to slide."""
        text_color = COLORS.TEXT_LIGHT if dark_bg else COLORS.TEXT_GRAY
        
        try:
            pos = self._get_container("footer")
            footer_box = slide.shapes.add_textbox(
                Cm(pos["left"]), Cm(pos["top"]),
                Cm(pos["width"]), Cm(pos["height"])
            )
            
            tf = footer_box.text_frame
            p = tf.paragraphs[0]
            p.text = self._get_footer_text()
            p.font.name = TYPOGRAPHY.BODY_FONT
            p.font.size = Pt(TYPOGRAPHY.FOOTER_SIZE)
            p.font.color.rgb = hex_to_rgb(text_color)
            
            if include_number:
                num_pos = self._get_container("slide_num")
                num_box = slide.shapes.add_textbox(
                    Cm(num_pos["left"]), Cm(num_pos["top"]),
                    Cm(num_pos["width"]), Cm(num_pos["height"])
                )
                tf = num_box.text_frame
                p = tf.paragraphs[0]
                p.text = str(self.slide_count)
                p.font.name = TYPOGRAPHY.BODY_FONT
                p.font.size = Pt(TYPOGRAPHY.FOOTER_SIZE)
                p.font.color.rgb = hex_to_rgb(text_color)
                p.alignment = PP_ALIGN.RIGHT
        except Exception as e:
            self.warnings.append(f"Failed to add footer: {e}")
    
    # -------------------------------------------------------------------------
    # SPEAKER NOTES METHOD
    # -------------------------------------------------------------------------
    
    def _add_speaker_notes(self, slide, notes_text: str):
        """Add speaker notes to a slide."""
        if not notes_text:
            return
        
        try:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes_text
        except Exception as e:
            self.warnings.append(f"Failed to add speaker notes: {e}")



    # Add auto-formatting helpers

    def _format_value(self, value: str, format_type: str = "auto") -> str:
        """Auto-format values for display."""
        if format_type == "auto":
            # Detect and format
            if value.startswith("$") or value.endswith("B") or value.endswith("M"):
                return value  # Already formatted
            try:
                num = float(value.replace(",", "").replace("%", ""))
                if "%" in value:
                    return f"{num:.1f}%"
                elif num >= 1_000_000_000:
                    return f"${num/1_000_000_000:.1f}B"
                elif num >= 1_000_000:
                    return f"${num/1_000_000:.1f}M"
                elif num >= 1_000:
                    return f"${num/1_000:.0f}K"
                else:
                    return f"${num:,.0f}"
            except:
                return value
        return value

    def _add_trend_indicator(self, value: str) -> str:
        """Add trend arrow if value contains growth rate."""
        if "%" in value:
            try:
                num = float(value.replace("%", "").replace("+", ""))
                if num > 0:
                    return f"↑ {value}"
                elif num < 0:
                    return f"↓ {value}"
            except:
                pass
        return value


        # =========================================================================
        # THEME APPLICATION METHODS
        # =========================================================================
        
        def apply_template_theme(self, template_path: str) -> bool:
            """
            Apply theme colors from a template presentation.
            
            This ensures the generated presentation uses the exact same
            color scheme as the corporate template.
            
            Args:
                template_path: Path to template.pptx with desired theme
                
            Returns:
                True if successful, False otherwise
            """
            try:
                from theme_extractor import ThemeApplier
                
                # Save current presentation to temp file
                import tempfile
                with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                    tmp_path = tmp.name
                    self.prs.save(tmp_path)
                
                # Apply theme from template
                ThemeApplier.copy_theme_from_template(
                    template_path,
                    tmp_path,
                    tmp_path
                )
                
                # Reload presentation with new theme
                self.prs = Presentation(tmp_path)
                
                # Clean up
                import os
                os.unlink(tmp_path)
                
                return True
                
            except Exception as e:
                self.warnings.append(f"Failed to apply template theme: {e}")
                return False
        
        @classmethod
        def from_template(
            cls,
            template_path: str,
            confidentiality: Confidentiality = Confidentiality.PUBLIC,
            year: Optional[int] = None,
            assets_path: str = "assets"
        ) -> 'ADIPresentation':
            """
            Create a new presentation based on a template.
            
            This copies the entire theme, slide masters, and layouts from
            the template, ensuring complete brand consistency.
            
            Args:
                template_path: Path to template.pptx
                confidentiality: Footer confidentiality level
                year: Year for copyright (defaults to current)
                assets_path: Path to brand assets
                
            Returns:
                New ADIPresentation instance with template theme applied
            """
            # Create presentation from template
            prs = Presentation(template_path)
            
            # Create instance
            instance = cls.__new__(cls)
            instance.prs = prs
            instance.confidentiality = confidentiality
            instance.year = year or datetime.now().year
            instance.assets = AssetManager(assets_path)
            instance.validator = BrandValidator()
            instance.slide_count = 0
            instance.warnings = []
            
            # Clear any existing slides from template (keep masters/layouts)
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
            
            return instance




    # =========================================================================
    # PUBLIC SLIDE METHODS
    # =========================================================================
    
    # =========================================================================
    # PUBLIC SLIDE METHODS (WITH SPEAKER NOTES SUPPORT)
    # =========================================================================
    
    def add_cover_slide(self, title: str, subtitle: str = None,
                        background_image: str = None,
                        use_default_bg: bool = True,
                        speaker_notes: str = None) -> 'ADIPresentation':
        """Add cover slide with AMP triangle and optional speaker notes."""
        self.slide_count += 1
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        # Background
        bg_added = False
        if background_image and Path(background_image).exists():
            bg_added = self._add_image_background(slide, Path(background_image))
        elif use_default_bg:
            bg_added = self._add_image_background(slide, self.assets.cover_background)
        
        if not bg_added:
            self._add_solid_background(slide, COLORS.DARK_BLUE)
        
        self._add_amp_overlay(slide)
        self._add_logo(slide, "cover", white=True)
        self._add_title(slide, title, "cover_title", dark_bg=True, size_pt=44)
        
        if subtitle:
            pos = self._get_container("cover_subtitle")
            try:
                sub_box = slide.shapes.add_textbox(
                    Cm(pos["left"]), Cm(pos["top"]),
                    Cm(pos["width"]), Cm(pos["height"])
                )
                tf = sub_box.text_frame
                p = tf.paragraphs[0]
                p.text = subtitle
                p.font.name = TYPOGRAPHY.BODY_FONT
                p.font.size = Pt(TYPOGRAPHY.SUBTITLE_SIZE)
                p.font.color.rgb = hex_to_rgb(COLORS.TEXT_LIGHT)
            except Exception as e:
                self.warnings.append(f"Failed to add subtitle: {e}")
        
        self._add_url(slide, "cover_url", dark_bg=True)
        
        # Add speaker notes
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)
        
        return self
    
    def add_section_slide(self, title: str,
                          slide_type: SectionSlideType = SectionSlideType.SECTION_TITLE,
                          background_image: str = None,
                          speaker_notes: str = None) -> 'ADIPresentation':
        """Add section divider slide with optional speaker notes."""
        self.slide_count += 1
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        if background_image and Path(background_image).exists():
            self._add_image_background(slide, Path(background_image))
        else:
            self._add_solid_background(slide, COLORS.DARK_BLUE)
        
        self._add_amp_overlay(slide)
        self._add_logo(slide, "cover", white=True)
        
        size = 48 if slide_type == SectionSlideType.KEY_MESSAGE else 44
        self._add_title(slide, title, "section_title", dark_bg=True, size_pt=size)
        
        self._add_footer(slide, dark_bg=True)
        
        # Add speaker notes
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)
        
        return self
    
    def add_content_slide(self, title: str,
                          content: List[Union[ContentItem, Tuple[str, int]]],
                          dark_background: bool = False,
                          speaker_notes: str = None) -> 'ADIPresentation':
        """Add content slide with optional speaker notes."""
        self.slide_count += 1
        
        items = []
        for c in content:
            if isinstance(c, ContentItem):
                items.append(c)
            elif isinstance(c, tuple) and len(c) >= 2:
                extra = c[2] if len(c) > 2 else False
                items.append(ContentItem(text=c[0], level=c[1], extra_space=extra))
            else:
                items.append(ContentItem(text=str(c), level=1))
        
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        bg_color = COLORS.DARK_BLUE if dark_background else COLORS.WHITE
        self._add_solid_background(slide, bg_color)
        self._add_logo(slide, "content", white=dark_background)
        self._add_title(slide, title, "content_title", dark_bg=dark_background)
        self._add_body_content(slide, items, "content_body", dark_bg=dark_background)
        self._add_footer(slide, dark_bg=dark_background)
        
        # Add speaker notes
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)
        
        return self
    
    def add_two_column_slide(self, title: str,
                              left_content: List[Union[ContentItem, Tuple[str, int]]],
                              right_content: List[Union[ContentItem, Tuple[str, int]]],
                              left_header: str = None,
                              right_header: str = None,
                              dark_background: bool = False,
                              speaker_notes: str = None) -> 'ADIPresentation':
        """Add two-column slide with optional headers and speaker notes."""
        self.slide_count += 1
        
        def convert_items(items):
            result = []
            for c in items:
                if isinstance(c, ContentItem):
                    result.append(c)
                elif isinstance(c, tuple) and len(c) >= 2:
                    extra = c[2] if len(c) > 2 else False
                    result.append(ContentItem(text=c[0], level=c[1], extra_space=extra))
                else:
                    result.append(ContentItem(text=str(c), level=1))
            return result
        
        # Prepend headers if provided
        left_items = convert_items(left_content)
        right_items = convert_items(right_content)
        
        if left_header:
            left_items.insert(0, ContentItem(text=left_header, level=1, extra_space=True))
        if right_header:
            right_items.insert(0, ContentItem(text=right_header, level=1, extra_space=True))
        
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        bg_color = COLORS.DARK_BLUE if dark_background else COLORS.WHITE
        self._add_solid_background(slide, bg_color)
        self._add_logo(slide, "content", white=dark_background)
        self._add_title(slide, title, "content_title", dark_bg=dark_background)
        self._add_body_content(slide, left_items, "two_col_left", dark_bg=dark_background)
        self._add_body_content(slide, right_items, "two_col_right", dark_bg=dark_background)
        self._add_footer(slide, dark_bg=dark_background)
        
        # Add speaker notes
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)
        
        return self
    
    def add_table_slide(self, title: str, table_data: TableData,
                        style: TableStyle = TableStyle.RECOMMENDED_14PT,
                        speaker_notes: str = None) -> 'ADIPresentation':
        """Add table slide with ADI styling and optional speaker notes."""
        self.slide_count += 1
        
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        self._add_solid_background(slide, COLORS.WHITE)
        self._add_logo(slide, "content", white=False)
        self._add_title(slide, title, "content_title", dark_bg=False)
        
        num_cols = len(table_data.headers)
        num_rows = len(table_data.rows) + 1
        
        pos = self._get_container("content_body")
        row_height = max(TABLE_CONFIG.MIN_ROW_HEIGHT_CM, min(12.0 / num_rows, 1.5))
        
        self.validator.validate_table_row_height(row_height)
        
        try:
            table_shape = slide.shapes.add_table(
                num_rows, num_cols,
                Cm(pos["left"]), Cm(pos["top"]),
                Cm(pos["width"]), Cm(row_height * num_rows)
            )
            table = table_shape.table
            
            col_width = int(Cm(pos["width"]) / num_cols)
            for col in table.columns:
                col.width = col_width
            
            # Style settings
            if style == TableStyle.DEFAULT_18PT:
                font_size = TABLE_CONFIG.DEFAULT_FONT_SIZE
            else:
                font_size = TABLE_CONFIG.RECOMMENDED_FONT_SIZE
            
            header_align = PP_ALIGN.CENTER if style in [TableStyle.CENTERED_14PT, TableStyle.TITLE_CENTER_14PT] else PP_ALIGN.LEFT
            cell_align = PP_ALIGN.CENTER if style == TableStyle.CENTERED_14PT else PP_ALIGN.LEFT
            
            # Headers
            for col_idx, header in enumerate(table_data.headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(TABLE_CONFIG.HEADER_BG_COLOR)
                
                p = cell.text_frame.paragraphs[0]
                p.font.name = TYPOGRAPHY.BODY_FONT
                p.font.size = Pt(font_size)
                p.font.bold = True
                p.alignment = header_align
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Data rows
            for row_idx, row in enumerate(table_data.rows):
                for col_idx, value in enumerate(row):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(value)
                    
                    # Highlight row if specified
                    if table_data.highlight_row is not None and row_idx == table_data.highlight_row:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = hex_to_rgb(COLORS.PRIMARY_BLUE)
                        p = cell.text_frame.paragraphs[0]
                        p.font.color.rgb = hex_to_rgb(COLORS.WHITE)
                        p.font.bold = True
                    elif row_idx % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = hex_to_rgb(TABLE_CONFIG.ALT_ROW_COLOR)
                    
                    p = cell.text_frame.paragraphs[0]
                    p.font.name = TYPOGRAPHY.BODY_FONT
                    p.font.size = Pt(font_size)
                    p.alignment = cell_align
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception as e:
            self.warnings.append(f"Failed to create table: {e}")
        
        self._add_footer(slide)
        
        # Add speaker notes
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)
        
        return self
    
    def add_chart_slide(self, title: str, chart_type: str,
                        categories: List[str], series: List[ChartSeries],
                        dark_background: bool = False,
                        chart_title: str = None,
                        speaker_notes: str = None) -> 'ADIPresentation':
        """Add chart slide with ADI colors and optional speaker notes."""
        self.slide_count += 1
        
        type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'area': XL_CHART_TYPE.AREA,
            'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
            'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
            'stacked_area': XL_CHART_TYPE.AREA_STACKED,
        }
        
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        bg_color = COLORS.DARK_BLUE if dark_background else COLORS.WHITE
        self._add_solid_background(slide, bg_color)
        self._add_logo(slide, "content", white=dark_background)
        self._add_title(slide, title, "content_title", dark_bg=dark_background)
        
        try:
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for s in series:
                chart_data.add_series(s.name, s.values)
            
            pos = self._get_container("content_body")
            xl_type = type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            chart_shape = slide.shapes.add_chart(
                xl_type,
                Cm(pos["left"]), Cm(pos["top"]),
                Cm(pos["width"]), Cm(pos["height"]),
                chart_data
            )
            
            chart = chart_shape.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            
            # Apply colors
            scheme = ChartColorScheme.DARK_BACKGROUND if dark_background else ChartColorScheme.LIGHT_BACKGROUND
            colors = COLORS.get_chart_colors(scheme)
            
            try:
                for i, s in enumerate(chart.series):
                    if i < len(colors):
                        s.format.fill.solid()
                        s.format.fill.fore_color.rgb = hex_to_rgb(colors[i])
            except Exception:
                pass
        except Exception as e:
            self.warnings.append(f"Failed to create chart: {e}")
        
        self._add_footer(slide, dark_bg=dark_background)
        
        # Add speaker notes
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)
        
        return self
    
    def add_closing_slide(self, speaker_notes: str = None) -> 'ADIPresentation':
        """Add closing slide with optional speaker notes."""
        self.slide_count += 1
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        if not self._add_image_background(slide, self.assets.tessellated_background):
            self._add_solid_background(slide, COLORS.DARK_BLUE)
        
        self._add_logo(slide, "closing", white=True)
        
        pos = self._get_container("closing_tagline")
        try:
            tag_box = slide.shapes.add_textbox(
                Cm(pos["left"]), Cm(pos["top"]),
                Cm(pos["width"]), Cm(pos["height"])
            )
            tf = tag_box.text_frame
            p = tf.paragraphs[0]
            p.text = CLOSING_TAGLINE
            p.font.name = TYPOGRAPHY.TITLE_FONT
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(COLORS.TEXT_LIGHT)
        except Exception as e:
            self.warnings.append(f"Failed to add tagline: {e}")
        
        self._add_url(slide, "closing_url", dark_bg=True)
        self._add_footer(slide, include_number=False, dark_bg=True)
        
        # Add speaker notes
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)
        
        return self
    
    # =========================================================================
    # NEW ADVANCED SLIDE METHODS
    # =========================================================================
    
    def add_metric_highlight_slide(
        self,
        title: str,
        metric_value: str,
        metric_label: str,
        supporting_text: str = None,
        dark_background: bool = True,
        speaker_notes: str = None
    ) -> 'ADIPresentation':
        """
        Add a slide with a single large metric for maximum impact.
        Ideal for key statistics like "$2.2B TAM" or "4-6 Hour Delay".
        """
        self.slide_count += 1
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        bg_color = COLORS.DARK_BLUE if dark_background else COLORS.WHITE
        text_color = COLORS.TEXT_LIGHT if dark_background else COLORS.TEXT_DARK
        
        self._add_solid_background(slide, bg_color)
        self._add_logo(slide, "content", white=dark_background)
        self._add_title(slide, title, "content_title", dark_bg=dark_background)
        
        # Large metric value
        try:
            metric_box = slide.shapes.add_textbox(
                Cm(2), Cm(6),
                Cm(29), Cm(5)
            )
            tf = metric_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric_value
            p.font.name = TYPOGRAPHY.TITLE_FONT
            p.font.size = Pt(72)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(COLORS.PRIMARY_BLUE if not dark_background else "#4D9AD4")
            p.alignment = PP_ALIGN.CENTER
            
            # Metric label
            label_box = slide.shapes.add_textbox(
                Cm(2), Cm(11),
                Cm(29), Cm(2)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric_label
            p.font.name = TYPOGRAPHY.BODY_FONT
            p.font.size = Pt(24)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            
            # Supporting text
            if supporting_text:
                support_box = slide.shapes.add_textbox(
                    Cm(4), Cm(13.5),
                    Cm(25), Cm(2)
                )
                tf = support_box.text_frame
                p = tf.paragraphs[0]
                p.text = supporting_text
                p.font.name = TYPOGRAPHY.BODY_FONT
                p.font.size = Pt(14)
                p.font.color.rgb = hex_to_rgb(COLORS.TEXT_GRAY if not dark_background else COLORS.TEXT_LIGHT)
                p.alignment = PP_ALIGN.CENTER
        except Exception as e:
            self.warnings.append(f"Failed to add metric: {e}")
        
        self._add_footer(slide, dark_bg=dark_background)
        
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)
        
        return self
    
    def add_comparison_slide(
        self,
        title: str,
        before_header: str,
        before_items: List[str],
        after_header: str,
        after_items: List[str],
        before_color: str = "#DC3545",  # Red
        after_color: str = "#28A745",   # Green
        speaker_notes: str = None
    ) -> 'ADIPresentation':
        """
        Add a before/after or problem/solution comparison slide with colored emphasis.
        """
        self.slide_count += 1
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        self._add_solid_background(slide, COLORS.WHITE)
        self._add_logo(slide, "content", white=False)
        self._add_title(slide, title, "content_title", dark_bg=False)
        
        # Left column (Before/Problem)
        try:
            left_header_box = slide.shapes.add_textbox(
                Cm(2), Cm(4),
                Cm(14), Cm(1.5)
            )
            tf = left_header_box.text_frame
            p = tf.paragraphs[0]
            p.text = before_header
            p.font.name = TYPOGRAPHY.TITLE_FONT
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(before_color)
            
            for i, item in enumerate(before_items[:5]):
                item_box = slide.shapes.add_textbox(
                    Cm(2), Cm(5.5 + i * 2),
                    Cm(14), Cm(1.8)
                )
                tf = item_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"✗ {item}"
                p.font.name = TYPOGRAPHY.BODY_FONT
                p.font.size = Pt(14)
                p.font.color.rgb = hex_to_rgb(COLORS.TEXT_DARK)
            
            # Right column (After/Solution)
            right_header_box = slide.shapes.add_textbox(
                Cm(17), Cm(4),
                Cm(14), Cm(1.5)
            )
            tf = right_header_box.text_frame
            p = tf.paragraphs[0]
            p.text = after_header
            p.font.name = TYPOGRAPHY.TITLE_FONT
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(after_color)
            
            for i, item in enumerate(after_items[:5]):
                item_box = slide.shapes.add_textbox(
                    Cm(17), Cm(5.5 + i * 2),
                    Cm(14), Cm(1.8)
                )
                tf = item_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"✓ {item}"
                p.font.name = TYPOGRAPHY.BODY_FONT
                p.font.size = Pt(14)
                p.font.color.rgb = hex_to_rgb(COLORS.TEXT_DARK)
        except Exception as e:
            self.warnings.append(f"Failed to add comparison content: {e}")
        
        self._add_footer(slide)
        
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)
        
        return self

    


    
    
    

    
    def add_blank_slide(self, dark_background: bool = False) -> 'ADIPresentation':
        """Add blank slide."""
        self.slide_count += 1
        slide = self.prs.slides.add_slide(self._blank_layout())
        
        bg_color = COLORS.DARK_BLUE if dark_background else COLORS.WHITE
        self._add_solid_background(slide, bg_color)
        self._add_logo(slide, "content", white=dark_background)
        self._add_footer(slide, dark_bg=dark_background)
        
        return self
    
    # =========================================================================
    # VALIDATION & SAVE
    # =========================================================================
    
    def validate(self) -> BrandValidator:
        """Run brand compliance validation."""
        self.validator.clear()
        self.validator.validate_slide_count(self.slide_count)
        return self.validator
    
    def save(self, filename: str, validate: bool = True) -> str:
        """Save presentation to file."""
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        
        if validate:
            self.validate()
            print(self.validator.get_report())
        
        if self.warnings:
            print("\n⚠️  Generation Warnings:")
            for w in self.warnings:
                print(f"   • {w}")
        
        try:
            self.prs.save(filename)
            print(f"\n✓ Saved: {filename}")
            print(f"  Slides: {self.slide_count}")
            print(f"  Confidentiality: {self.confidentiality.value}")
        except Exception as e:
            raise IOError(f"Failed to save presentation: {e}")
        
        return filename


# =============================================================================
# FACTORY FUNCTION
# =============================================================================

def create_presentation(
    confidentiality: str = 'public',
    year: Optional[int] = None,
    assets_path: str = "assets"
) -> ADIPresentation:
    """Factory function to create ADI presentation."""
    conf_map = {
        'public': Confidentiality.PUBLIC,
        'confidential': Confidentiality.CONFIDENTIAL,
        'internal_only': Confidentiality.INTERNAL_ONLY,
    }
    return ADIPresentation(
        confidentiality=conf_map.get(confidentiality.lower(), Confidentiality.PUBLIC),
        year=year,
        assets_path=assets_path
    )
